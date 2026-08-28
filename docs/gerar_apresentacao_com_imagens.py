"""
Script para geração da Apresentação de Defesa Técnica da Avaliação 3 (Versão 16:9 Dark Executive)
com layout otimizado contendo ESPAÇO DEDICADO PARA INSERÇÃO DE IMAGENS / SCREENSHOTS.

Equipe 01 - Disciplina: Técnicas de Hyperautomation - Prof. Moisés Levy.
Estrutura adaptada para 4 integrantes, com divisão equilibrada, notas de orador completas
e cards/dropzones específicos para screenshots do projeto.
"""

import pptx
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pathlib import Path

# --- PALETA DE CORES PROFISSIONAL HYPERAUTOMATION / LG ---
COLOR_BG_DARK = RGBColor(11, 19, 43)        # #0B132B - Navy Profundo
COLOR_BG_CARD = RGBColor(28, 37, 65)        # #1C2541 - Card Azul Escuro
COLOR_BG_CARD_LIGHT = RGBColor(37, 50, 85)  # #253255 - Card Destaque
COLOR_BG_IMAGE_ZONE = RGBColor(18, 26, 48)  # #121A30 - Fundo para Dropzone de Imagens
COLOR_ACCENT_RED = RGBColor(165, 0, 52)     # #A50034 - LG Ruby Red
COLOR_ACCENT_CYAN = RGBColor(0, 180, 216)   # #00B4D8 - Electric Cyan
COLOR_ACCENT_GOLD = RGBColor(255, 209, 102) # #FFD166 - Gold / Warning
COLOR_ACCENT_GREEN = RGBColor(6, 214, 160)  # #06D6A0 - Emerald Green
COLOR_TEXT_WHITE = RGBColor(255, 255, 255)  # Branco Puro
COLOR_TEXT_MUTED = RGBColor(200, 205, 220)  # Cinza Claro
COLOR_BORDER = RGBColor(58, 80, 107)        # Borda sutil
COLOR_BORDER_DASHED = RGBColor(0, 180, 216) # Borda da Dropzone de Imagem


def criar_apresentacao_com_imagens(caminho_saida: str):
    prs = pptx.Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    def aplicar_fundo_padrao(slide, titulo="", categoria="", integrante=""):
        # Fundo Geral
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
        bg.fill.solid()
        bg.fill.fore_color.rgb = COLOR_BG_DARK
        bg.line.fill.background()

        if titulo:
            # Header Container
            header_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.733), Inches(1.1))
            tf = header_box.text_frame
            tf.word_wrap = True
            tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0

            # Categoria / Tag superior
            if categoria:
                p_cat = tf.paragraphs[0]
                p_cat.text = categoria.upper()
                p_cat.font.size = Pt(10)
                p_cat.font.bold = True
                p_cat.font.color.rgb = COLOR_ACCENT_CYAN
                p_cat.space_after = Pt(2)
                p_tit = tf.add_paragraph()
            else:
                p_tit = tf.paragraphs[0]

            # Título Principal
            p_tit.text = titulo
            p_tit.font.size = Pt(21)
            p_tit.font.bold = True
            p_tit.font.color.rgb = COLOR_TEXT_WHITE

            # Badge do Integrante Responsável (Canto Superior Direito)
            if integrante:
                badge_w, badge_h = Inches(3.4), Inches(0.45)
                badge_x, badge_y = Inches(13.333 - 0.8 - 3.4), Inches(0.45)
                badge_bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, badge_x, badge_y, badge_w, badge_h)
                badge_bg.fill.solid()
                badge_bg.fill.fore_color.rgb = COLOR_ACCENT_RED
                badge_bg.line.fill.background()
                
                tf_b = badge_bg.text_frame
                tf_b.vertical_anchor = MSO_ANCHOR.MIDDLE
                p_b = tf_b.paragraphs[0]
                p_b.text = f"👤 {integrante}"
                p_b.font.size = Pt(10)
                p_b.font.bold = True
                p_b.font.color.rgb = COLOR_TEXT_WHITE
                p_b.alignment = PP_ALIGN.CENTER

        # Linha Divisória Superior
        div = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.5), Inches(11.733), Inches(0.03))
        div.fill.solid()
        div.fill.fore_color.rgb = COLOR_BORDER
        div.line.fill.background()

        # Rodapé Institucional
        foot_box = slide.shapes.add_textbox(Inches(0.8), Inches(7.0), Inches(11.733), Inches(0.35))
        tf_f = foot_box.text_frame
        tf_f.margin_left = tf_f.margin_top = tf_f.margin_right = tf_f.margin_bottom = 0
        p_f = tf_f.paragraphs[0]
        p_f.text = "Técnicas de Hyperautomation (Prof. Moisés Levy) • LG Electronics do Brasil • Equipe 01 • Turma T02"
        p_f.font.size = Pt(9)
        p_f.font.color.rgb = COLOR_BORDER

    def adicionar_card_texto(slide, x, y, w, h, titulo="", bg_color=COLOR_BG_CARD, border_color=COLOR_BORDER):
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
        card.fill.solid()
        card.fill.fore_color.rgb = bg_color
        card.line.color.rgb = border_color
        card.line.width = Pt(1.2)

        tf = card.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.2)
        tf.margin_top = Inches(0.2)
        tf.margin_right = Inches(0.2)
        tf.margin_bottom = Inches(0.2)

        if titulo:
            p_tit = tf.paragraphs[0]
            p_tit.text = titulo
            p_tit.font.size = Pt(13)
            p_tit.font.bold = True
            p_tit.font.color.rgb = COLOR_ACCENT_CYAN
            p_tit.space_after = Pt(6)

        return card

    def adicionar_espaco_imagem(slide, x, y, w, h, titulo_area="📷 ÁREA DE IMAGEM / SCREENSHOT", dica="Insira aqui o screenshot ou diagrama correspondente"):
        """Cria um espaço reservado destacado especificamente para colar/inserir imagens."""
        # Container de fundo
        card_img = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
        card_img.fill.solid()
        card_img.fill.fore_color.rgb = COLOR_BG_IMAGE_ZONE
        card_img.line.color.rgb = COLOR_BORDER_DASHED
        card_img.line.width = Pt(1.5)

        tf = card_img.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = Inches(0.2)
        tf.margin_right = Inches(0.2)

        p1 = tf.paragraphs[0]
        p1.text = titulo_area
        p1.font.size = Pt(12)
        p1.font.bold = True
        p1.font.color.rgb = COLOR_ACCENT_GOLD
        p1.alignment = PP_ALIGN.CENTER
        p1.space_after = Pt(6)

        p2 = tf.add_paragraph()
        p2.text = dica
        p2.font.size = Pt(9.5)
        p2.font.color.rgb = COLOR_TEXT_MUTED
        p2.alignment = PP_ALIGN.CENTER

        return card_img

    # =========================================================================
    # SLIDE 1: CAPA DO PROJETO (👤 Integrante 1)
    # =========================================================================
    slide1 = prs.slides.add_slide(blank_layout)
    bg1 = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    bg1.fill.solid()
    bg1.fill.fore_color.rgb = COLOR_BG_DARK
    bg1.line.fill.background()

    # Bloco de Título
    t_box1 = slide1.shapes.add_textbox(Inches(0.8), Inches(1.0), Inches(7.2), Inches(3.5))
    tf1 = t_box1.text_frame
    tf1.word_wrap = True

    p0 = tf1.paragraphs[0]
    p0.text = "AVALIAÇÃO 03 — DEFESA TÉCNICA E APRESENTAÇÃO DE PROJETO"
    p0.font.size = Pt(11)
    p0.font.bold = True
    p0.font.color.rgb = COLOR_ACCENT_RED
    p0.space_after = Pt(8)

    p1 = tf1.add_paragraph()
    p1.text = "Automação Inteligente do Processo de\nSeleção de Fornecedores Industriais"
    p1.font.size = Pt(27)
    p1.font.bold = True
    p1.font.color.rgb = COLOR_TEXT_WHITE
    p1.space_after = Pt(10)

    p2 = tf1.add_paragraph()
    p2.text = "Pipeline ponta a ponta com RPA (Playwright), Algoritmo de Decisão Multicritério (MCDA), Containerização Docker e CI/CD"
    p2.font.size = Pt(12)
    p2.font.color.rgb = COLOR_ACCENT_CYAN

    # Espaço para Logo / Banner / Imagem de Capa no lado direito
    adicionar_espaco_imagem(
        slide1, Inches(8.3), Inches(1.0), Inches(4.233), Inches(3.4),
        "📷 [ ESPAÇO PARA LOGO / BANNER DA LG / CAPA ]",
        "Insira aqui o logotipo da LG Electronics / AX Academy ou imagem ilustrativa de Hyperautomation."
    )

    # Card da Equipe
    card_eq = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(4.7), Inches(11.733), Inches(2.0))
    card_eq.fill.solid()
    card_eq.fill.fore_color.rgb = COLOR_BG_CARD
    card_eq.line.color.rgb = COLOR_BORDER
    tf_eq = card_eq.text_frame
    tf_eq.word_wrap = True
    tf_eq.margin_left = Inches(0.3)
    tf_eq.margin_top = Inches(0.18)

    peq1 = tf_eq.paragraphs[0]
    peq1.text = "EQUIPE 01 — DIVISÃO DA DEFESA EM 4 INTEGRANTES (Turma T02 • Prof. Moisés Levy):"
    peq1.font.size = Pt(11)
    peq1.font.bold = True
    peq1.font.color.rgb = COLOR_ACCENT_GOLD
    peq1.space_after = Pt(4)

    peq2 = tf_eq.add_paragraph()
    peq2.text = "• Integrante 1 (Líder / Negócio): Problema AS-IS na LG, Solução Proposta TO-BE e Governança\n" \
                "• Integrante 2 (Automação / Dados): Pipeline em 6 Etapas, Arquitetura Técnica e Cenário Normal (Fornecedores A, B, C)\n" \
                "• Integrante 3 (Qualidade / MCDA): Gestão de Exceções (Fornecedor D), Resiliência e Suíte Pytest (51 Testes)\n" \
                "• Integrante 4 (DevOps / Governança): Docker, Docker Compose, CI/CD GitHub Actions, GHCR e Auditoria SOX"
    peq2.font.size = Pt(10.5)
    peq2.font.color.rgb = COLOR_TEXT_WHITE

    slide1.notes_slide.notes_text_frame.text = (
        "NOTAS DO ORADOR (👤 INTEGRANTE 1):\n"
        "- Cumprimentar o professor Moisés Levy e a banca avaliadora.\n"
        "- Apresentar a Equipe 01 e introduzir o tema: Automação Inteligente do Processo de Seleção de Fornecedores para a LG Electronics do Brasil.\n"
        "- Informar que a apresentação está estruturada de forma lógica seguindo a esteira oficial: Problema -> Solução -> Automação -> Testes -> Docker -> CI/CD -> GHCR -> Execução -> Resultado.\n"
        "- Explicar que a defesa está dividida de maneira equilibrada entre os 4 integrantes."
    )

    # =========================================================================
    # SLIDE 2: O PROBLEMA AS-IS (👤 Integrante 1)
    # =========================================================================
    slide2 = prs.slides.add_slide(blank_layout)
    aplicar_fundo_padrao(slide2, "O Problema de Negócio: Processo Manual AS-IS na LG", "1. Contexto & Dores", "Integrante 1 (Líder / Negócio)")

    c1 = adicionar_card_texto(slide2, Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.9), "FLUXO MANUAL & IMPACTOS (AS-IS)")
    tf = c1.text_frame
    p = tf.add_paragraph()
    p.text = "1. Disparo Manual de RFQs por E-mail:\n" \
             "   Compradores solicitam cotações sem padronização.\n\n" \
             "2. Recepção Despadronizada:\n" \
             "   Propostas chegam em formatos dispersos (.xlsx, .csv).\n\n" \
             "3. Transcrição Manual para Planilhas Locais:\n" \
             "   Comprador digita preços, prazos e impostos no Excel.\n\n" \
             "🚨 Impactos e Riscos Críticos:\n" \
             "• Erro Humano: 12% de divergências em valores e prazos.\n" \
             "• Alto Lead Time: Ciclo de compras levava de 5 a 10 dias úteis.\n" \
             "• Não-Conformidade SOX: Falta de logs e auditoria digital."
    p.font.size = Pt(10.5)
    p.font.color.rgb = COLOR_TEXT_MUTED

    adicionar_espaco_imagem(
        slide2, Inches(6.8), Inches(1.8), Inches(5.733), Inches(4.9),
        "📷 [ DIAGRAMA DO FLUXO MANUAL AS-IS ]",
        "Cole aqui o fluxograma do processo manual anterior (E-mail -> Leitura Manual -> Digitação no Excel -> Risco de Erro)."
    )

    slide2.notes_slide.notes_text_frame.text = (
        "NOTAS DO ORADOR (👤 INTEGRANTE 1):\n"
        "- Explicar o cenário AS-IS da LG Electronics em Manaus.\n"
        "- Destacar que no modelo anterior, o comprador dependia de transcrição manual de propostas recebidas por e-mail em formatos dispersos.\n"
        "- Enfatizar as três dores críticas: risco de erro em cálculos de TCO, demora de até 10 dias no ciclo de compra e vulnerabilidade para auditorias SOX."
    )

    # =========================================================================
    # SLIDE 3: A SOLUÇÃO TO-BE (👤 Integrante 1)
    # =========================================================================
    slide3 = prs.slides.add_slide(blank_layout)
    aplicar_fundo_padrao(slide3, "A Solução Proposta: Hyperautomation de Ponta a Ponta", "2. Visão TO-BE", "Integrante 1 (Líder / Negócio)")

    c1 = adicionar_card_texto(slide3, Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.9), "PILAREs DA SOLUÇÃO (TO-BE)")
    tf = c1.text_frame
    p = tf.add_paragraph()
    p.text = "1. Ingestão Automatizada & RPA Web:\n" \
             "• Varredura autônoma de diretórios (.xlsx e .csv).\n" \
             "• Robô Playwright consulta status no portal corporativo.\n\n" \
             "2. Validação & Motor de Decisão MCDA:\n" \
             "• Barramento instantâneo de dados negativos e bloqueios.\n" \
             "• Algoritmo de scoring relativo com pesos oficiais de negócio.\n\n" \
             "3. Governança, CI/CD & Auditoria SOX:\n" \
             "• Planilha homologada oficial e JSON de auditoria digital.\n" \
             "• Redução do tempo de processamento de 8h para 25 segundos."
    p.font.size = Pt(10.5)
    p.font.color.rgb = COLOR_TEXT_MUTED

    adicionar_espaco_imagem(
        slide3, Inches(6.8), Inches(1.8), Inches(5.733), Inches(4.9),
        "📷 [ ARQUITETURA CONCEITUAL TO-BE ]",
        "Cole aqui o diagrama da solução de Hyperautomation (Coleta RPA -> Validação em Camadas -> Algoritmo MCDA -> Auditoria SOX)."
    )

    slide3.notes_slide.notes_text_frame.text = (
        "NOTAS DO ORADOR (👤 INTEGRANTE 1):\n"
        "- Apresentar a visão TO-BE dividida em Ingestão/RPA, Validação/Decisão e Governança/Auditoria.\n"
        "- Enfatizar a redução de 95% do tempo de processamento e conformidade de 100% com auditorias.\n"
        "- Passar a palavra para o Integrante 2."
    )

    # =========================================================================
    # SLIDE 4: FLUXO DO PIPELINE EM 6 ETAPAS (👤 Integrante 2)
    # =========================================================================
    slide4 = prs.slides.add_slide(blank_layout)
    aplicar_fundo_padrao(slide4, "Fluxo da Solução: Pipeline Modular em 6 Etapas", "3. Arquitetura do Pipeline", "Integrante 2 (Automação / Dados)")

    c1 = adicionar_card_texto(slide4, Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.9), "ESTRUTURA DAS 6 ETAPAS")
    tf = c1.text_frame
    p = tf.add_paragraph()
    p.text = "• ETAPA 1 (Coleta): Varredura de arquivos e Playwright RPA no portal web.\n\n" \
             "• ETAPA 2 (Leitura): Ingestão multiformato (Excel/CSV) e carga de critérios.\n\n" \
             "• ETAPA 3 (Validação): Barramento cadastral e numérico (Data Filtering).\n\n" \
             "• ETAPA 4 (Consolidação): Tipagem estrita e unificação em DataFrame.\n\n" \
             "• ETAPA 5 (Ranking MCDA): Scoring relativo ponderado por critérios.\n\n" \
             "• ETAPA 6 (Resultado & Auditoria): Preenchimento do template oficial e log SOX."
    p.font.size = Pt(10.5)
    p.font.color.rgb = COLOR_TEXT_MUTED

    adicionar_espaco_imagem(
        slide4, Inches(6.8), Inches(1.8), Inches(5.733), Inches(4.9),
        "📷 [ FLUXOGRAMA DO PIPELINE / 6 ETAPAS ]",
        "Cole aqui o fluxograma visual do pipeline (Coleta -> Leitura -> Validação -> Consolidação -> Ranking -> Resultado)."
    )

    slide4.notes_slide.notes_text_frame.text = (
        "NOTAS DO ORADOR (👤 INTEGRANTE 2):\n"
        "- Assumir a apresentação explicando a arquitetura em 6 etapas da pipeline.\n"
        "- Explicar como cada módulo possui responsabilidade bem definida e desacoplada, facilitando testes e manutenção."
    )

    # =========================================================================
    # SLIDE 5: ARQUITETURA E TECNOLOGIAS (👤 Integrante 2)
    # =========================================================================
    slide5 = prs.slides.add_slide(blank_layout)
    aplicar_fundo_padrao(slide5, "Stack Tecnológico e Integração dos Componentes", "4. Engenharia de Software", "Integrante 2 (Automação / Dados)")

    c1 = adicionar_card_texto(slide5, Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.9), "TECNOLOGIAS UTILIZADAS")
    tf = c1.text_frame
    p = tf.add_paragraph()
    p.text = "🐍 Python 3.12 (Core Engine):\n" \
             "• Orquestrador modular com tipagem estrita.\n\n" \
             "🌐 Playwright (Chromium RPA Headless):\n" \
             "• Automação web moderna e rápida sem dependência de WebDriver.\n\n" \
             "📊 Pandas & OpenPyXL:\n" \
             "• Manipulação matricial vetorial e geração de relatórios Excel.\n\n" \
             "🧪 Pytest & Flake8:\n" \
             "• 51 testes automatizados e 0 erros de linting PEP8.\n\n" \
             "🐳 Docker, Docker Compose & GHCR:\n" \
             "• Timezone America/Manaus e publicação oficial no ghcr.io."
    p.font.size = Pt(10.5)
    p.font.color.rgb = COLOR_TEXT_MUTED

    adicionar_espaco_imagem(
        slide5, Inches(6.8), Inches(1.8), Inches(5.733), Inches(4.9),
        "📷 [ DIAGRAMA DE COMPONENTES E ARQUITETURA ]",
        "Cole aqui o diagrama com os logos e conexões das tecnologias (Python, Playwright, Pandas, Pytest, Docker, GitHub Actions, GHCR)."
    )

    slide5.notes_slide.notes_text_frame.text = (
        "NOTAS DO ORADOR (👤 INTEGRANTE 2):\n"
        "- Justificar as escolhas técnicas: Playwright moderno para RPA web headless, Pandas para operações vetoriais de MCDA, Pytest para testes e Docker Compose com timezone America/Manaus."
    )

    # =========================================================================
    # SLIDE 6: CENÁRIO NORMAL / CAMINHO FELIZ (👤 Integrante 2)
    # =========================================================================
    slide6 = prs.slides.add_slide(blank_layout)
    aplicar_fundo_padrao(slide6, "Cenário de Execução Normal: Fornecedores A, B e C", "5. Demonstração Cenário Normal", "Integrante 2 (Automação / Dados)")

    c1 = adicionar_card_texto(slide6, Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.9), "ENTRADA & HOMOLOGAÇÃO NORMAL")
    tf = c1.text_frame
    p = tf.add_paragraph()
    p.text = "1. Propostas Comerciais Processadas:\n" \
             "• Fornecedor A (.xlsx): Custo R$ 100 | Prazo: 5d | Cap: 500 | Qual: 95\n" \
             "• Fornecedor B (.csv):  Custo R$ 90  | Prazo: 8d | Cap: 700 | Qual: 90\n" \
             "• Fornecedor C (.xlsx): Custo R$ 110 | Prazo: 4d | Cap: 400 | Qual: 98\n\n" \
             "2. Consulta Cadastral Playwright: Status 'Ativo' validado.\n\n" \
             "🏆 Classificação Oficial Gerada:\n" \
             "• 1º Lugar: Fornecedor B (Score 92.98) -> Menor Custo (R$ 90)\n" \
             "• 2º Lugar: Fornecedor A (Score 89.28) -> Proposta Equilibrada\n" \
             "• 3º Lugar: Fornecedor C (Score 88.08) -> Melhor Prazo/Qualidade"
    p.font.size = Pt(10.5)
    p.font.color.rgb = COLOR_TEXT_MUTED

    adicionar_espaco_imagem(
        slide6, Inches(6.8), Inches(1.8), Inches(5.733), Inches(4.9),
        "📷 [ SCREENSHOT DO TERMINAL / EXECUÇÃO NORMAL ]",
        "Cole aqui o print da execução do robô no terminal mostrando os 4 arquivos localizados e os 3 fornecedores válidos aprovados."
    )

    slide6.notes_slide.notes_text_frame.text = (
        "NOTAS DO ORADOR (👤 INTEGRANTE 2):\n"
        "- Demonstrar o caminho feliz da esteira.\n"
        "- Mostrar que os 3 fornecedores válidos foram processados e o Fornecedor B venceu pelo peso preponderante de Custo (40%).\n"
        "- Passar a palavra para o Integrante 3."
    )

    # =========================================================================
    # SLIDE 7: TRATAMENTO DE EXCEÇÕES E RESILIÊNCIA (👤 Integrante 3)
    # =========================================================================
    slide7 = prs.slides.add_slide(blank_layout)
    aplicar_fundo_padrao(slide7, "Tratamento de Exceções, Resiliência e Barramentos", "6. Gestão de Falhas & Exceções", "Integrante 3 (Qualidade / MCDA)")

    c1 = adicionar_card_texto(slide7, Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.9), "REGRAS DE BARRAMENTO & RESILIÊNCIA")
    tf = c1.text_frame
    p = tf.add_paragraph()
    p.text = "🚨 Desclassificação do Fornecedor D:\n" \
             "1. Bloqueio Cadastral Web: Status 'Bloqueado' no portal HTML.\n" \
             "2. Dados Negativos: Custo (-50), Prazo (-2) e Capacidade (-100).\n" \
             "🛡️ Data Filtering Gateway: Segregado para não distorcer as notas dos demais concorrentes (Nota 0.00).\n\n" \
             "🔄 Resiliência Multi-Camada de Rede:\n" \
             "• Nível 1: Playwright HTTP (localhost:8000)\n" \
             "• Nível 2: Playwright Local (file://) em nova aba limpa\n" \
             "• Nível 3: Requests + BeautifulSoup\n" \
             "• Nível 4: Leitura direta em disco (0% de parada na fábrica)."
    p.font.size = Pt(10.5)
    p.font.color.rgb = COLOR_TEXT_MUTED

    adicionar_espaco_imagem(
        slide7, Inches(6.8), Inches(1.8), Inches(5.733), Inches(4.9),
        "📷 [ SCREENSHOT DO PORTAL WEB / FORNECEDOR D BLOQUEADO ]",
        "Cole aqui o print do portal web simulado mostrando o Fornecedor D 'Bloqueado' ou o log de rejeição no terminal."
    )

    slide7.notes_slide.notes_text_frame.text = (
        "NOTAS DO ORADOR (👤 INTEGRANTE 3):\n"
        "- Assumir a apresentação respondendo: 'Como o robô trata erros?'.\n"
        "- Explicar o caso do Fornecedor D (bloqueado e dados negativos) e como o Data Filtering Gateway evita distorcer as notas dos demais.\n"
        "- Demonstrar a resiliência de rede em 4 camadas."
    )

    # =========================================================================
    # SLIDE 8: TESTES AUTOMATIZADOS PYTEST (👤 Integrante 3)
    # =========================================================================
    slide8 = prs.slides.add_slide(blank_layout)
    aplicar_fundo_padrao(slide8, "Suíte de Testes Automatizados com Pytest (51 Testes)", "7. Qualidade & Testes", "Integrante 3 (Qualidade / MCDA)")

    c1 = adicionar_card_texto(slide8, Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.9), "COBERTURA & ESTRUTURA DE TESTES")
    tf = c1.text_frame
    p = tf.add_paragraph()
    p.text = "🧪 test_coleta.py (24 testes):\n" \
             "• Filtros de arquivos de lock temporários (~$*).\n" \
             "• Playwright RPA com Chromium real e mocks de erro de porta.\n\n" \
             "🧪 test_leitura.py (25 testes):\n" \
             "• Ingestão de XLSX e CSV com delimitadores ';' e ','.\n" \
             "• Tolerância a encodings (utf-8, latin-1) e acentuações.\n" \
             "• Fallback automático de critérios e modelo de ranking.\n\n" \
             "🧪 test_integration.py (2 testes):\n" \
             "• Fluxo integrado ponta a ponta e execução via main.py.\n\n" \
             "📊 Resultado: 51 Testes Aprovados (100% Passed ✅) em ~35s."
    p.font.size = Pt(10.5)
    p.font.color.rgb = COLOR_TEXT_MUTED

    adicionar_espaco_imagem(
        slide8, Inches(6.8), Inches(1.8), Inches(5.733), Inches(4.9),
        "📷 [ SCREENSHOT DO PYTEST - 51 PASSED / COVERAGE ]",
        "Cole aqui o print do terminal executando 'python -m pytest -v' com todos os 51 testes passando em verde."
    )

    slide8.notes_slide.notes_text_frame.text = (
        "NOTAS DO ORADOR (👤 INTEGRANTE 3):\n"
        "- Apresentar a suíte com 51 testes automatizados cobrindo todos os ramos de negócio e exceção com 100% de sucesso.\n"
        "- Destacar a integração no CI/CD."
    )

    # =========================================================================
    # SLIDE 9: ALGORITMO MCDA PONDERADO (👤 Integrante 3)
    # =========================================================================
    slide9 = prs.slides.add_slide(blank_layout)
    aplicar_fundo_padrao(slide9, "Motor de Decisão Multicritério: Algoritmo MCDA", "8. Modelo Matemático", "Integrante 3 (Qualidade / MCDA)")

    c1 = adicionar_card_texto(slide9, Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.9), "PESOS DE NEGÓCIO & EQUAÇÕES")
    tf = c1.text_frame
    p = tf.add_paragraph()
    p.text = "📋 Matriz de Critérios (criterios_ranking.xlsx):\n" \
             "• Custo (40% / 0.40) | Direção: Menor é Melhor 🔻\n" \
             "• Prazo (25% / 0.25) | Direção: Menor é Melhor 🔻\n" \
             "• Capacidade (20% / 0.20) | Direção: Maior é Melhor 🔺\n" \
             "• Qualidade (15% / 0.15) | Direção: Maior é Melhor 🔺\n" \
             "⚖️ Soma dos Pesos = 1.00 (100%)\n\n" \
             "📐 Normalização Relativa ao Benchmark de Mercado:\n" \
             "• Menor é Melhor: (Menor Valor / Valor da Proposta) * 100\n" \
             "• Maior é Melhor: (Valor da Proposta / Maior Valor) * 100\n\n" \
             "🎯 Score Global: Soma ponderada dos 4 critérios."
    p.font.size = Pt(10.5)
    p.font.color.rgb = COLOR_TEXT_MUTED

    adicionar_espaco_imagem(
        slide9, Inches(6.8), Inches(1.8), Inches(5.733), Inches(4.9),
        "📷 [ TABELA DE CRITÉRIOS / MEMÓRIA DE CÁLCULO ]",
        "Cole aqui o print da planilha criterios_ranking.xlsx ou da memória de cálculo dos scores parciais."
    )

    slide9.notes_slide.notes_text_frame.text = (
        "NOTAS DO ORADOR (👤 INTEGRANTE 3):\n"
        "- Explicar as fórmulas matemáticas do MCDA e o benchmark de mercado.\n"
        "- Passar a palavra para o Integrante 4."
    )

    # =========================================================================
    # SLIDE 10: DOCKER & DOCKER COMPOSE (👤 Integrante 4)
    # =========================================================================
    slide10 = prs.slides.add_slide(blank_layout)
    aplicar_fundo_padrao(slide10, "Containerização com Docker & Docker Compose", "9. Infraestrutura & Containers", "Integrante 4 (DevOps / Governança)")

    c1 = adicionar_card_texto(slide10, Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.9), "DOCKERFILE & COMPOSE MULTI-SERVIÇO")
    tf = c1.text_frame
    p = tf.add_paragraph()
    p.text = "🐳 Imagem Base: python:3.12-slim otimizada.\n\n" \
             "🕒 Timezone Oficial: TZ=America/Manaus\n" \
             "• Configurado no SO (tzdata) e nas variáveis de ambiente para rastreabilidade de Manaus.\n\n" \
             "🌐 Playwright Chromium Integrado no Build.\n\n" \
             "📂 Volumes Persistentes:\n" \
             "• ./output -> /app/output (Planilha gerada)\n" \
             "• ./logs   -> /app/logs   (Auditoria JSON e logs)\n\n" \
             "🚀 Docker Compose Multi-Serviço:\n" \
             "• Serviço 'web-panel' (porta 8000) + Serviço 'robot'."
    p.font.size = Pt(10.5)
    p.font.color.rgb = COLOR_TEXT_MUTED

    adicionar_espaco_imagem(
        slide10, Inches(6.8), Inches(1.8), Inches(5.733), Inches(4.9),
        "📷 [ SCREENSHOT DOCKER / DOCKER COMPOSE ]",
        "Cole aqui o print do terminal executando 'docker compose up --build' ou a visualização no Docker Desktop."
    )

    slide10.notes_slide.notes_text_frame.text = (
        "NOTAS DO ORADOR (👤 INTEGRANTE 4):\n"
        "- Assumir a apresentação abordando a infraestrutura de Docker e Docker Compose.\n"
        "- Explicar os volumes persistentes, as variáveis de ambiente e o timezone America/Manaus."
    )

    # =========================================================================
    # SLIDE 11: CI/CD GITHUB ACTIONS & GHCR (👤 Integrante 4)
    # =========================================================================
    slide11 = prs.slides.add_slide(blank_layout)
    aplicar_fundo_padrao(slide11, "Pipeline de CI/CD (GitHub Actions) & Publicação GHCR", "10. DevOps & Integração Contínua", "Integrante 4 (DevOps / Governança)")

    c1 = adicionar_card_texto(slide11, Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.9), "ESTEIRA CI/CD (.github/workflows/ci.yml)")
    tf = c1.text_frame
    p = tf.add_paragraph()
    p.text = "🔄 Gatilhos: Push e Pull Request em main, develop e feature/**.\n\n" \
             "⚙️ Job 1: Lint, Testes e Cobertura:\n" \
             "• Setup Python 3.12 com cache pip e Chromium do Playwright.\n" \
             "• Linting estrito Flake8 (0 erros).\n" \
             "• 51 testes Pytest com relatório de cobertura XML.\n\n" \
             "⚙️ Job 2: Build & Publicação no GHCR:\n" \
             "• Imagem Docker publicada no GitHub Container Registry (ghcr.io).\n\n" \
             "💻 Execução em Produção:\n" \
             "   docker pull ghcr.io/sannyer3232/hyperautomation-av3-equipe1:latest"
    p.font.size = Pt(10.5)
    p.font.color.rgb = COLOR_TEXT_MUTED

    adicionar_espaco_imagem(
        slide11, Inches(6.8), Inches(1.8), Inches(5.733), Inches(4.9),
        "📷 [ SCREENSHOT GITHUB ACTIONS & GHCR ]",
        "Cole aqui o print da esteira de CI com checks verdes no GitHub Actions e a imagem publicada na aba Packages (GHCR)."
    )

    slide11.notes_slide.notes_text_frame.text = (
        "NOTAS DO ORADOR (👤 INTEGRANTE 4):\n"
        "- Explicar a esteira automatizada de CI/CD no GitHub Actions e a publicação no GitHub Container Registry (GHCR).\n"
        "- Mostrar os comandos docker pull e docker run."
    )

    # =========================================================================
    # SLIDE 12: LOGS E AUDITORIA SOX (👤 Integrante 4)
    # =========================================================================
    slide12 = prs.slides.add_slide(blank_layout)
    aplicar_fundo_padrao(slide12, "Rastreabilidade Digital, Logs e Auditoria SOX", "11. Governança & Compliance", "Integrante 4 (DevOps / Governança)")

    c1 = adicionar_card_texto(slide12, Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.9), "GOVERNANÇA & ARTEFATOS FINAIS")
    tf = c1.text_frame
    p = tf.add_paragraph()
    p.text = "🏛️ Trilha de Auditoria SOX (logs/auditoria.json):\n" \
             "• Carimbo de data/hora ISO 8601 (TZ Manaus).\n" \
             "• Total de propostas recebidas (4) e aprovadas (3).\n" \
             "• Justificativa da desclassificação do Fornecedor D.\n" \
             "• Memória de cálculo matemática completa dos scores.\n\n" \
             "📋 Logs Operacionais (logs/execucao.log):\n" \
             "• 4 Níveis: [INFO], [WARNING], [ERROR], [CRITICAL].\n\n" \
             "📊 Planilha Homologada (output/ranking_final.xlsx):\n" \
             "• 1º Lugar: Fornecedor B | 2º A | 3º C | Desclassificado: D."
    p.font.size = Pt(10.5)
    p.font.color.rgb = COLOR_TEXT_MUTED

    adicionar_espaco_imagem(
        slide12, Inches(6.8), Inches(1.8), Inches(5.733), Inches(4.9),
        "📷 [ SCREENSHOT DO RANKING_FINAL.XLSX & AUDITORIA.JSON ]",
        "Cole aqui o print da planilha oficial gerada no Excel e do arquivo auditoria.json formatado."
    )

    slide12.notes_slide.notes_text_frame.text = (
        "NOTAS DO ORADOR (👤 INTEGRANTE 4):\n"
        "- Explicar a diferença entre o log operacional para suporte e o snapshot JSON de auditoria SOX para compliance da LG."
    )

    # =========================================================================
    # SLIDE 13: CONCLUSÃO E VALOR ENTREGUE (👤 Integrante 4)
    # =========================================================================
    slide13 = prs.slides.add_slide(blank_layout)
    aplicar_fundo_padrao(slide13, "Conclusão e Entrega de Valor para o Negócio", "12. Fechamento & ROI", "Integrante 4 (DevOps / Governança)")

    c1 = adicionar_card_texto(slide13, Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.9), "RESUMO EXECUTIVO & ROI")
    tf = c1.text_frame
    p = tf.add_paragraph()
    p.text = "1. O que foi desenvolvido?\n" \
             "• Solução autônoma de Hyperautomation com RPA, MCDA, Docker e CI/CD.\n\n" \
             "2. Qual problema resolveu?\n" \
             "• Eliminação de 100% da digitação manual de propostas.\n" \
             "• Barramento de fornecedores irregulares.\n\n" \
             "3. Qual resultado foi obtido?\n" \
             "• Redução do tempo de ciclo de 8h para 25 segundos.\n\n" \
             "4. Qual valor entrega para a LG?\n" \
             "• Decisões objetivas de compras e conformidade SOX total."
    p.font.size = Pt(10.5)
    p.font.color.rgb = COLOR_TEXT_MUTED

    adicionar_espaco_imagem(
        slide13, Inches(6.8), Inches(1.8), Inches(5.733), Inches(4.9),
        "📷 [ GRÁFICO DE GANHOS / QUADRO DE RESULTADOS ]",
        "Cole aqui um gráfico de barras comparativo (Tempo Manual 8h vs Robô 25s) ou imagem de encerramento dos resultados."
    )

    slide13.notes_slide.notes_text_frame.text = (
        "NOTAS DO ORADOR (👤 INTEGRANTE 4):\n"
        "- Fazer o fechamento executivo respondendo às 5 perguntas obrigatórias de conclusão.\n"
        "- Anunciar o início da demonstração prática ao vivo."
    )

    # =========================================================================
    # SLIDE 14: ROTEIRO DA DEMONSTRAÇÃO PRÁTICA (👥 4 Integrantes)
    # =========================================================================
    slide14 = prs.slides.add_slide(blank_layout)
    aplicar_fundo_padrao(slide14, "Roteiro da Demonstração Prática (7 Etapas Obrigatórias)", "Demonstração ao Vivo", "4 Integrantes (Demonstração)")

    c1 = adicionar_card_texto(slide14, Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.9), "ROTEIRO PASSO A PASSO DA DEMO")
    tf = c1.text_frame
    p = tf.add_paragraph()
    p.text = "👤 Integrante 1 — ETAPA 1: Código & Pastas no VS Code\n" \
             "👤 Integrante 2 — ETAPA 2: Execução 'python src/main.py'\n" \
             "👤 Integrante 3 — ETAPA 3: Testes 'python -m pytest -v'\n" \
             "👤 Integrante 4 — ETAPA 4: GitHub Actions no Navegador\n" \
             "👤 Integrante 4 — ETAPA 5: Imagem Publicada no GHCR\n" \
             "👤 Integrante 4 — ETAPA 6: Container 'docker compose up'\n" \
             "👥 Todos — ETAPA 7: Exibição de ranking_final.xlsx e auditoria.json"
    p.font.size = Pt(10.5)
    p.font.color.rgb = COLOR_TEXT_MUTED

    adicionar_espaco_imagem(
        slide14, Inches(6.8), Inches(1.8), Inches(5.733), Inches(4.9),
        "📷 [ CHECKLIST VISUAL DA DEMO AO VIVO ]",
        "Cole aqui o checklist ou esquema das 7 etapas da banca (Código -> Execução -> Pytest -> Actions -> GHCR -> Docker -> Resultados)."
    )

    slide14.notes_slide.notes_text_frame.text = (
        "NOTAS DO ORADOR (👥 TODOS OS 4 INTEGRANTES):\n"
        "- Seguir exatamente este roteiro cronometrado na demonstração ao vivo para a banca examinadora."
    )

    # =========================================================================
    # SLIDE 15: ENCERRAMENTO E DÚVIDAS (👥 4 Integrantes)
    # =========================================================================
    slide15 = prs.slides.add_slide(blank_layout)
    aplicar_fundo_padrao(slide15, "Defesa Técnica Individual: Perguntas da Banca & Dúvidas", "Perguntas & Respostas", "4 Integrantes")

    c1 = adicionar_card_texto(slide15, Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.9), "PERGUNTAS FREQUENTES DA BANCA")
    tf = c1.text_frame
    p = tf.add_paragraph()
    p.text = "1. 'Por que Playwright?' -> Velocidade, headless nativo e estabilidade.\n\n" \
             "2. 'Por que o Fornecedor D foi rejeitado?' -> Bloqueado na web e dados negativos.\n\n" \
             "3. 'Como funciona o cálculo MCDA?' -> Normalização de mercado + pesos oficiais.\n\n" \
             "4. 'Como se recupera de falha de rede?' -> Fallback Playwright HTTP -> file:// -> Disco.\n\n" \
             "5. 'Onde estão os volumes no Docker?' -> ./output e ./logs persistidos no host."
    p.font.size = Pt(10.0)
    p.font.color.rgb = COLOR_TEXT_MUTED

    c2 = adicionar_card_texto(slide15, Inches(6.8), Inches(1.8), Inches(5.733), Inches(4.9), "OBRIGADO! DÚVIDAS DA BANCA?", bg_color=COLOR_BG_CARD_LIGHT)
    tf = c2.text_frame
    p = tf.add_paragraph()
    p.text = "LG ELECTRONICS DO BRASIL\n" \
             "Automação Inteligente do Processo de Seleção de Fornecedores\n\n" \
             "👨‍💻 Equipe 01 — Turma T02 (2026)\n" \
             "• Sannyer Carvalho & Integrantes da Equipe\n\n" \
             "🎓 Disciplina: Técnicas de Hyperautomation\n" \
             "👨‍🏫 Professor Responsável: Prof. Moisés Levy\n\n" \
             "🔗 Repositório: github.com/Sannyer3232/hyperautomation-av3-equipe1\n" \
             "📦 Imagem GHCR: ghcr.io/sannyer3232/hyperautomation-av3-equipe1:latest\n\n" \
             "Estamos abertos para as perguntas da banca avaliadora!"
    p.font.size = Pt(10.5)
    p.font.color.rgb = COLOR_TEXT_WHITE

    slide15.notes_slide.notes_text_frame.text = (
        "NOTAS DO ORADOR (👥 TODOS OS 4 INTEGRANTES):\n"
        "- Agradecer a atenção do professor Moisés Levy e da banca.\n"
        "- Abrir para a sessão de perguntas individuais e coletivas da banca examinadora."
    )

    prs.save(caminho_saida)
    print(f"Sucesso! Apresentação com espaços de imagem gerada em: {caminho_saida}")


if __name__ == "__main__":
    caminho = Path("docs") / "Apresentacao_Defesa_Hyperautomation_LG_Equipe01.pptx"
    criar_apresentacao_com_imagens(str(caminho))
