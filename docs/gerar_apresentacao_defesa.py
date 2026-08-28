"""
Script para geração automática dos slides da Apresentação de Defesa Técnica da Avaliação 3.
Projeto: Automação Inteligente do Processo de Seleção de Fornecedores (LG Electronics).
Equipe 01 - Disciplina: Técnicas de Hyperautomation - Prof. Moisés Levy.
Estrutura adaptada para divisão equilibrada entre 4 integrantes.
"""

import pptx
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pathlib import Path

# --- PALETA DE CORES PROFISSIONAL HYPERAUTOMATION / LG ---
COLOR_BG_DARK = RGBColor(11, 19, 43)        # #0B132B - Navy Profundo
COLOR_BG_CARD = RGBColor(28, 37, 65)        # #1C2541 - Card Azul Escuro
COLOR_BG_CARD_LIGHT = RGBColor(37, 50, 85)  # #253255 - Card Destaque
COLOR_ACCENT_RED = RGBColor(165, 0, 52)     # #A50034 - LG Ruby Red
COLOR_ACCENT_CYAN = RGBColor(0, 180, 216)   # #00B4D8 - Electric Cyan
COLOR_ACCENT_GOLD = RGBColor(255, 209, 102) # #FFD166 - Gold / Warning
COLOR_ACCENT_GREEN = RGBColor(6, 214, 160)  # #06D6A0 - Emerald Green
COLOR_TEXT_WHITE = RGBColor(255, 255, 255)  # Branco Puro
COLOR_TEXT_MUTED = RGBColor(200, 205, 220)  # Cinza Claro
COLOR_BORDER = RGBColor(58, 80, 107)        # Borda sutil


def criar_apresentacao_completa(caminho_saida: str):
    prs = pptx.Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    def aplicar_fundo_padrao(slide, titulo="", categoria="", integrante=""):
        # Fundo Geral
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
        bg.fill.solid()
        bg.fill.fore_color.rgb = COLOR_BG_DARK
        bg.line.fill.background()

        if titulo:
            # Header Container
            header_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.733), Inches(1.1))
            tf = header_box.text_frame
            tf.word_wrap = True
            tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0

            # Categoria / Tag superior
            if categoria:
                p_cat = tf.paragraphs[0]
                p_cat.text = categoria.upper()
                p_cat.font.size = Pt(10)
                p_cat.font.bold = True
                p_cat.font.color.rgb = COLOR_ACCENT_CYAN
                p_cat.space_after = Pt(2)
                p_tit = tf.add_paragraph()
            else:
                p_tit = tf.paragraphs[0]

            # Título Principal
            p_tit.text = titulo
            p_tit.font.size = Pt(22)
            p_tit.font.bold = True
            p_tit.font.color.rgb = COLOR_TEXT_WHITE

            # Badge do Integrante Responsável (Canto Superior Direito)
            if integrante:
                badge_w, badge_h = Inches(3.2), Inches(0.45)
                badge_x, badge_y = Inches(13.333 - 0.8 - 3.2), Inches(0.45)
                badge_bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, badge_x, badge_y, badge_w, badge_h)
                badge_bg.fill.solid()
                badge_bg.fill.fore_color.rgb = COLOR_ACCENT_RED
                badge_bg.line.fill.background()
                
                tf_b = badge_bg.text_frame
                tf_b.vertical_anchor = MSO_ANCHOR.MIDDLE
                p_b = tf_b.paragraphs[0]
                p_b.text = f"👤 {integrante}"
                p_b.font.size = Pt(11)
                p_b.font.bold = True
                p_b.font.color.rgb = COLOR_TEXT_WHITE
                p_b.alignment = PP_ALIGN.CENTER

        # Linha Divisória Superior
        div = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.5), Inches(11.733), Inches(0.03))
        div.fill.solid()
        div.fill.fore_color.rgb = COLOR_BORDER
        div.line.fill.background()

        # Rodapé Institucional
        foot_box = slide.shapes.add_textbox(Inches(0.8), Inches(7.0), Inches(11.733), Inches(0.35))
        tf_f = foot_box.text_frame
        tf_f.margin_left = tf_f.margin_top = tf_f.margin_right = tf_f.margin_bottom = 0
        p_f = tf_f.paragraphs[0]
        p_f.text = "Técnicas de Hyperautomation (Prof. Moisés Levy) • LG Electronics do Brasil • Equipe 01 • Turma T02"
        p_f.font.size = Pt(9)
        p_f.font.color.rgb = COLOR_BORDER

    def adicionar_card(slide, x, y, w, h, titulo="", bg_color=COLOR_BG_CARD, border_color=COLOR_BORDER):
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
        card.fill.solid()
        card.fill.fore_color.rgb = bg_color
        card.line.color.rgb = border_color
        card.line.width = Pt(1.2)

        tf = card.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.2)
        tf.margin_top = Inches(0.2)
        tf.margin_right = Inches(0.2)
        tf.margin_bottom = Inches(0.2)

        if titulo:
            p_tit = tf.paragraphs[0]
            p_tit.text = titulo
            p_tit.font.size = Pt(14)
            p_tit.font.bold = True
            p_tit.font.color.rgb = COLOR_ACCENT_CYAN
            p_tit.space_after = Pt(8)

        return card

    # =========================================================================
    # SLIDE 1: CAPA DO PROJETO (👤 Integrante 1)
    # =========================================================================
    slide1 = prs.slides.add_slide(blank_layout)
    bg1 = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    bg1.fill.solid()
    bg1.fill.fore_color.rgb = COLOR_BG_DARK
    bg1.line.fill.background()

    # Tag Superior
    t_box1 = slide1.shapes.add_textbox(Inches(1.0), Inches(1.2), Inches(11.333), Inches(4.5))
    tf1 = t_box1.text_frame
    tf1.word_wrap = True

    p0 = tf1.paragraphs[0]
    p0.text = "AVALIAÇÃO 03 — DEFESA TÉCNICA E APRESENTAÇÃO DE PROJETO"
    p0.font.size = Pt(12)
    p0.font.bold = True
    p0.font.color.rgb = COLOR_ACCENT_RED
    p0.space_after = Pt(10)

    p1 = tf1.add_paragraph()
    p1.text = "Automação Inteligente do Processo de\nSeleção de Fornecedores Industriais"
    p1.font.size = Pt(30)
    p1.font.bold = True
    p1.font.color.rgb = COLOR_TEXT_WHITE
    p1.space_after = Pt(12)

    p2 = tf1.add_paragraph()
    p2.text = "Pipeline ponta a ponta com RPA (Playwright), Algoritmo de Decisão Multicritério (MCDA), Containerização Docker e CI/CD"
    p2.font.size = Pt(14)
    p2.font.color.rgb = COLOR_ACCENT_CYAN
    p2.space_after = Pt(28)

    # Card da Equipe
    card_eq = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(4.7), Inches(11.333), Inches(2.0))
    card_eq.fill.solid()
    card_eq.fill.fore_color.rgb = COLOR_BG_CARD
    card_eq.line.color.rgb = COLOR_BORDER
    tf_eq = card_eq.text_frame
    tf_eq.word_wrap = True
    tf_eq.margin_left = Inches(0.3)
    tf_eq.margin_top = Inches(0.2)

    peq1 = tf_eq.paragraphs[0]
    peq1.text = "EQUIPE 01 — DIVISÃO DA DEFESA EM 4 INTEGRANTES:"
    peq1.font.size = Pt(11)
    peq1.font.bold = True
    peq1.font.color.rgb = COLOR_ACCENT_GOLD
    peq1.space_after = Pt(6)

    peq2 = tf_eq.add_paragraph()
    peq2.text = "• Integrante 1 (Líder / Negócio): Problema AS-IS, Solução Proposta TO-BE e Governança\n" \
                "• Integrante 2 (Automação / Dados): Pipeline de 6 Etapas, Arquitetura Técnica e Cenário Normal (A, B, C)\n" \
                "• Integrante 3 (Qualidade / MCDA): Tratamento de Exceções (Fornecedor D), Resiliência e Testes Pytest\n" \
                "• Integrante 4 (DevOps / Governança): Docker, Docker Compose, CI/CD GitHub Actions, GHCR e Auditoria SOX"
    peq2.font.size = Pt(11)
    peq2.font.color.rgb = COLOR_TEXT_WHITE

    slide1.notes_slide.notes_text_frame.text = (
        "NOTAS DO ORADOR (👤 INTEGRANTE 1):\n"
        "- Cumprimentar o professor Moisés Levy e a banca avaliadora.\n"
        "- Apresentar a Equipe 01 e introduzir o tema: Automação Inteligente do Processo de Seleção de Fornecedores para a LG Electronics do Brasil.\n"
        "- Informar que a apresentação está estruturada de forma lógica seguindo a esteira oficial: Problema -> Solução -> Automação -> Testes -> Docker -> CI/CD -> GHCR -> Execução -> Resultado.\n"
        "- Explicar que a defesa está dividida de maneira equilibrada entre os 4 integrantes."
    )

    # =========================================================================
    # SLIDE 2: O PROBLEMA AS-IS (👤 Integrante 1)
    # =========================================================================
    slide2 = prs.slides.add_slide(blank_layout)
    aplicar_fundo_padrao(slide2, "O Problema de Negócio: Processo Manual AS-IS na LG", "1. Contexto & Dores", "Integrante 1")

    c1 = adicionar_card(slide2, Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.9), "FLUXO OPERACIONAL MANUAL (AS-IS)")
    tf_c1 = c1.text_frame
    p = tf_c1.add_paragraph()
    p.text = "1. Disparo Manual de RFQs:\n" \
             "   Compradores enviam cotações por e-mail para dezenas de fornecedores.\n\n" \
             "2. Recepção Despadronizada:\n" \
             "   Propostas chegam em formatos heterogêneos (.xlsx, .csv, corpos de e-mail).\n\n" \
             "3. Transcrição e Digitação no Excel:\n" \
             "   Comprador digita manualmente preços, prazos e tributos em planilhas locais.\n\n" \
             "4. Consulta Fragmentada de Compliance:\n" \
             "   Checagem manual de certidões e status cadastral em portais externos.\n\n" \
             "5. Cálculo Manual de TCO e Parecer:\n" \
             "   Fórmulas manuais suscetíveis a erros antes da aprovação da gerência."
    p.font.size = Pt(11)
    p.font.color.rgb = COLOR_TEXT_MUTED

    c2 = adicionar_card(slide2, Inches(6.8), Inches(1.8), Inches(5.733), Inches(4.9), "IMPACTOS & VULNERABILIDADES", bg_color=COLOR_BG_CARD_LIGHT)
    tf_c2 = c2.text_frame
    p = tf_c2.add_paragraph()
    p.text = "🚨 Erros Humanos e Fórmulas Inconsistentes:\n" \
             "• 12% de divergências em transcrição de valores e prazos.\n" \
             "• Risco de homologação de fornecedores bloqueados ou com dados negativos.\n\n" \
             "⏱️ Alto Lead Time de Suprimentos:\n" \
             "• Ciclo de compras levava de 5 a 10 dias úteis (~8h de esforço manual direto).\n\n" \
             "🔍 Baixa Rastreabilidade e Não-Conformidade SOX:\n" \
             "• Falta de logs estruturados e histórico transparente das decisões.\n" \
             "• Dificuldade em justificar auditorias corporativas e critérios ponderados."
    p.font.size = Pt(11)
    p.font.color.rgb = COLOR_TEXT_WHITE

    slide2.notes_slide.notes_text_frame.text = (
        "NOTAS DO ORADOR (👤 INTEGRANTE 1):\n"
        "- Explicar o cenário AS-IS da LG Electronics em Manaus.\n"
        "- Destacar que no modelo anterior, o comprador dependia de transcrição manual de propostas recebidas por e-mail em formatos dispersos (Excel e CSV).\n"
        "- Enfatizar as três dores críticas: risco de erro em cálculos de TCO, demora de até 10 dias no ciclo de compra e vulnerabilidade para auditorias SOX por falta de rastreabilidade digital.\n"
        "- Concluir dizendo que essa ineficiência exigia uma solução de Hyperautomation de ponta a ponta."
    )

    # =========================================================================
    # SLIDE 3: A SOLUÇÃO TO-BE (👤 Integrante 1)
    # =========================================================================
    slide3 = prs.slides.add_slide(blank_layout)
    aplicar_fundo_padrao(slide3, "A Solução Proposta: Hyperautomation de Ponta a Ponta", "2. Visão TO-BE", "Integrante 1")

    c1 = adicionar_card(slide3, Inches(0.8), Inches(1.8), Inches(3.6), Inches(4.9), "1. INGESTÃO & RPA")
    tf = c1.text_frame
    p = tf.add_paragraph()
    p.text = "• Varredura automática de arquivos (.xlsx e .csv).\n" \
             "• Robô Playwright acessa portal web simulado.\n" \
             "• Extração de status cadastrais (Ativo/Bloqueado).\n" \
             "• Fallback inteligente em camadas."
    p.font.size = Pt(11)
    p.font.color.rgb = COLOR_TEXT_MUTED

    c2 = adicionar_card(slide3, Inches(4.8), Inches(1.8), Inches(3.7), Inches(4.9), "2. VALIDAÇÃO & DECISÃO")
    tf = c2.text_frame
    p = tf.add_paragraph()
    p.text = "• Gateway de validação cadastral e numérica.\n" \
             "• Rejeição instantânea de dados negativos e bloqueios.\n" \
             "• Algoritmo MCDA ponderado com normalização relativa.\n" \
             "• Pesos dinâmicos oficiais."
    p.font.size = Pt(11)
    p.font.color.rgb = COLOR_TEXT_MUTED

    c3 = adicionar_card(slide3, Inches(8.9), Inches(1.8), Inches(3.633), Inches(4.9), "3. GOVERNANÇA & AUDITORIA")
    tf = c3.text_frame
    p = tf.add_paragraph()
    p.text = "• Preenchimento do template oficial homologado.\n" \
             "• Geração do ranking_final.xlsx.\n" \
             "• Trilha de auditoria SOX (auditoria.json).\n" \
             "• Logs contínuos com 4 níveis de severidade."
    p.font.size = Pt(11)
    p.font.color.rgb = COLOR_TEXT_MUTED

    slide3.notes_slide.notes_text_frame.text = (
        "NOTAS DO ORADOR (👤 INTEGRANTE 1):\n"
        "- Apresentar a arquitetura conceitual TO-BE dividida em 3 pilares: Ingestão/RPA, Validação/Decisão e Governança/Auditoria.\n"
        "- Explicar que a solução não apenas lê planilhas, mas integra navegador automatizado (Playwright) para checar o status dos fornecedores no portal, aplica regras rígidas de barramento para propostas inválidas, calcula o score ponderado de mercado e gera a planilha oficial preenchida.\n"
        "- Passar a palavra para o Integrante 2, que explicará o fluxo detalhado das 6 etapas e a arquitetura técnica."
    )

    # =========================================================================
    # SLIDE 4: FLUXO DO PIPELINE EM 6 ETAPAS (👤 Integrante 2)
    # =========================================================================
    slide4 = prs.slides.add_slide(blank_layout)
    aplicar_fundo_padrao(slide4, "Fluxo da Solução: Pipeline Modular em 6 Etapas", "3. Arquitetura do Pipeline", "Integrante 2")

    largura_etapa = Inches(1.8)
    espacamento = Inches(0.18)
    
    etapas = [
        ("ETAPA 1\nColeta", "• Varredura .xlsx/.csv\n• Playwright RPA\n• Portal Web HTTP\n• Fallback local", COLOR_ACCENT_CYAN),
        ("ETAPA 2\nLeitura", "• Leitura pandas\n• Tolerância delimitadores\n• Normalização colunas\n• Carga de critérios", COLOR_ACCENT_CYAN),
        ("ETAPA 3\nValidação", "• Checagem portal web\n• Barramento negativo\n• Rejeição Fornec. D\n• Data Filtering", COLOR_ACCENT_GOLD),
        ("ETAPA 4\nConsolidação", "• Unificação DataFrames\n• Tipagem estrita\n• Isolamento rejeitadas\n• Preparação MCDA", COLOR_ACCENT_GOLD),
        ("ETAPA 5\nRanking", "• Scoring relativo\n• Pesos ponderados\n• Menor custo/prazo\n• Maior capac./qualid.", COLOR_ACCENT_GREEN),
        ("ETAPA 6\nResultado", "• Template oficial\n• ranking_final.xlsx\n• Auditoria SOX JSON\n• Logs detalhados", COLOR_ACCENT_GREEN)
    ]

    for i, (tit_e, desc_e, cor_e) in enumerate(etapas):
        x_pos = Inches(0.8) + i * (largura_etapa + espacamento)
        c = adicionar_card(slide4, x_pos, Inches(1.8), largura_etapa, Inches(4.9), tit_e, bg_color=COLOR_BG_CARD, border_color=cor_e)
        tf = c.text_frame
        p = tf.add_paragraph()
        p.text = desc_e
        p.font.size = Pt(10)
        p.font.color.rgb = COLOR_TEXT_WHITE

    slide4.notes_slide.notes_text_frame.text = (
        "NOTAS DO ORADOR (👤 INTEGRANTE 2):\n"
        "- Assumir a apresentação explicando a modularidade da pipeline em 6 etapas sequenciais.\n"
        "- Explicar que cada etapa possui responsabilidade única e isolada: Coleta de arquivos e web -> Leitura e padronização multiformato -> Validação técnica e de compliance -> Consolidação de dados válidos -> Cálculo matemático MCDA -> Geração do resultado homologado e auditoria SOX.\n"
        "- Destacar que essa separação em módulos garante facilidade de manutenção, testes unitários isolados e desacoplamento do código."
    )

    # =========================================================================
    # SLIDE 5: ARQUITETURA E TECNOLOGIAS (👤 Integrante 2)
    # =========================================================================
    slide5 = prs.slides.add_slide(blank_layout)
    aplicar_fundo_padrao(slide5, "Stack Tecnológico e Integração dos Componentes", "4. Engenharia de Software", "Integrante 2")

    c1 = adicionar_card(slide5, Inches(0.8), Inches(1.8), Inches(3.6), Inches(4.9), "LINGUAGEM & ROBÓTICA")
    tf = c1.text_frame
    p = tf.add_paragraph()
    p.text = "🐍 Python 3.12 (Core):\n" \
             "• Orquestrador modular ponta a ponta.\n\n" \
             "🌐 Playwright (Chromium RPA):\n" \
             "• Automação web headless rápida e moderna.\n" \
             "• Extração de tabelas e status cadastrais.\n\n" \
             "📊 Pandas & OpenPyXL:\n" \
             "• Processamento de Excel, CSV e exportação final."
    p.font.size = Pt(11)
    p.font.color.rgb = COLOR_TEXT_MUTED

    c2 = adicionar_card(slide5, Inches(4.8), Inches(1.8), Inches(3.7), Inches(4.9), "QUALIDADE & TESTES")
    tf = c2.text_frame
    p = tf.add_paragraph()
    p.text = "🧪 Pytest (v8.0+):\n" \
             "• 51 testes automatizados (unitários e integração).\n\n" \
             "📈 Pytest-Cov:\n" \
             "• Relatórios de cobertura de código.\n\n" \
             "🧹 Flake8 & Black:\n" \
             "• Padronização estrita PEP8 e 0 erros de linting."
    p.font.size = Pt(11)
    p.font.color.rgb = COLOR_TEXT_MUTED

    c3 = adicionar_card(slide5, Inches(8.9), Inches(1.8), Inches(3.633), Inches(4.9), "INFRA & DEVOPS")
    tf = c3.text_frame
    p = tf.add_paragraph()
    p.text = "🐳 Docker & Docker Compose:\n" \
             "• Containers isolados (Robô + Web Server).\n" \
             "• Timezone oficial: America/Manaus.\n\n" \
             "⚙️ GitHub Actions & GHCR:\n" \
             "• Esteira de CI/CD completa.\n" \
             "• Publicação no GitHub Container Registry."
    p.font.size = Pt(11)
    p.font.color.rgb = COLOR_TEXT_MUTED

    slide5.notes_slide.notes_text_frame.text = (
        "NOTAS DO ORADOR (👤 INTEGRANTE 2):\n"
        "- Detalhar por que cada tecnologia foi escolhida.\n"
        "- Explicar o uso do Playwright moderno em vez de BotCity ou Selenium: maior velocidade, execução headless nativa e resiliência com seletores DOM.\n"
        "- Explicar o uso de Pandas para a manipulação matricial e cálculo matemático do MCDA.\n"
        "- Mencionar a suíte robusta com Pytest, Flake8 e Docker Compose que prepara a aplicação para rodar em produção."
    )

    # =========================================================================
    # SLIDE 6: CENÁRIO NORMAL / CAMINHO FELIZ (👤 Integrante 2)
    # =========================================================================
    slide6 = prs.slides.add_slide(blank_layout)
    aplicar_fundo_padrao(slide6, "Cenário de Execução Normal: Fornecedores A, B e C", "5. Demonstração Cenário Normal", "Integrante 2")

    c1 = adicionar_card(slide6, Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.9), "ENTRADA & PROCESSAMENTO NORMAL")
    tf = c1.text_frame
    p = tf.add_paragraph()
    p.text = "1. Recepção das Propostas Comerciais:\n" \
             "• Fornecedor A (.xlsx): Custo R$ 100 | Prazo: 5d | Cap: 500 | Qual: 95\n" \
             "• Fornecedor B (.csv):  Custo R$ 90  | Prazo: 8d | Cap: 700 | Qual: 90\n" \
             "• Fornecedor C (.xlsx): Custo R$ 110 | Prazo: 4d | Cap: 400 | Qual: 98\n\n" \
             "2. Consulta Cadastral Web via Playwright:\n" \
             "• Fornecedores A, B e C identificados como 'Ativo'.\n\n" \
             "3. Validação Técnica Aprovada:\n" \
             "• Todos os valores positivos e dentro dos limites operacionais.\n\n" \
             "4. Consolidação no DataFrame Válido:\n" \
             "• 100% dos registros válidos encaminhados para o cálculo MCDA."
    p.font.size = Pt(11)
    p.font.color.rgb = COLOR_TEXT_MUTED

    c2 = adicionar_card(slide6, Inches(6.8), Inches(1.8), Inches(5.733), Inches(4.9), "RESULTADO DO CAMINHO FELIZ", bg_color=COLOR_BG_CARD_LIGHT)
    tf = c2.text_frame
    p = tf.add_paragraph()
    p.text = "🏆 Homologação Automática e Ranking Gerado:\n\n" \
             "• 1º Lugar: Fornecedor B (Score Final: 92.98)\n" \
             "  -> Menor custo de mercado (R$ 90) e alta capacidade produtiva (700 un).\n\n" \
             "• 2º Lugar: Fornecedor A (Score Final: 89.28)\n" \
             "  -> Proposta equilibrada entre prazo (5d) e qualidade (95%).\n\n" \
             "• 3º Lugar: Fornecedor C (Score Final: 88.08)\n" \
             "  -> Melhor prazo (4d) e maior qualidade (98%), porém custo mais alto (R$ 110).\n\n" \
             "✅ Resultado exportado em output/ranking_final.xlsx com formatação oficial da LG."
    p.font.size = Pt(11)
    p.font.color.rgb = COLOR_TEXT_WHITE

    slide6.notes_slide.notes_text_frame.text = (
        "NOTAS DO ORADOR (👤 INTEGRANTE 2):\n"
        "- Demonstrar o caminho feliz da esteira.\n"
        "- Explicar que os 3 fornecedores válidos (A, B e C) passaram por leitura multiformato (.xlsx e .csv), tiveram seus status 'Ativo' validados no portal web pelo Playwright e foram consolidados sem inconsistências.\n"
        "- Mostrar o ranking gerado onde o Fornecedor B venceu por liderar no critério de maior peso (Custo 40%).\n"
        "- Passar a palavra para o Integrante 3, que explicará o tratamento de exceções (Fornecedor D), resiliência e os testes automatizados."
    )

    # =========================================================================
    # SLIDE 7: TRATAMENTO DE EXCEÇÕES E RESILIÊNCIA (👤 Integrante 3)
    # =========================================================================
    slide7 = prs.slides.add_slide(blank_layout)
    aplicar_fundo_padrao(slide7, "Tratamento de Exceções, Resiliência e Barramentos", "6. Gestão de Falhas & Exceções", "Integrante 3")

    c1 = adicionar_card(slide7, Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.9), "CASO 1: DESCLASSIFICAÇÃO FORNECEDOR D", border_color=COLOR_ACCENT_RED)
    tf = c1.text_frame
    p = tf.add_paragraph()
    p.text = "🚨 Identificação Dupla de Não-Conformidade:\n" \
             "1. Camada Cadastral Web:\n" \
             "   • Status 'Bloqueado' detectado na tabela HTML do portal.\n" \
             "2. Camada de Integridade Numérica:\n" \
             "   • Custo negativo (-50), Prazo negativo (-2), Capacidade negativa (-100).\n\n" \
             "🛡️ Padrão Data Filtering Gateway:\n" \
             "• A proposta D é imediatamente segregada em 'propostas_rejeitadas'.\n" \
             "• Não entra no cálculo matemático do MCDA (evita distorção dos scores dos demais).\n" \
             "• Figura no ranking final como 'Desclassificado' (Nota 0.00) com o motivo registrado na coluna Observacao."
    p.font.size = Pt(11)
    p.font.color.rgb = COLOR_TEXT_MUTED

    c2 = adicionar_card(slide7, Inches(6.8), Inches(1.8), Inches(5.733), Inches(4.9), "CASO 2: RESILIÊNCIA DE REDE & FALLBACK", border_color=COLOR_ACCENT_CYAN)
    tf = c2.text_frame
    p = tf.add_paragraph()
    p.text = "🔄 Arquitetura Multi-Camada de Contingência Web:\n\n" \
             "• Nível 1 (Principal): Playwright via HTTP (http://localhost:8000/...)\n" \
             "  -> Se o servidor web responder normalmente, executa o scraping.\n\n" \
             "• Nível 2 (Contingência Playwright): Playwright via file://\n" \
             "  -> Se a porta 8000 estiver fechada ou timeout, abre nova aba limpa no arquivo HTML local.\n\n" \
             "• Nível 3 (Fallback HTTP): Requests + BeautifulSoup\n" \
             "  -> Se o Chromium falhar por falta de dependências gráficas.\n\n" \
             "• Nível 4 (Fallback Local): Leitura direta em disco\n" \
             "  -> Garantia de 0% de interrupção operacional na fábrica."
    p.font.size = Pt(11)
    p.font.color.rgb = COLOR_TEXT_WHITE

    slide7.notes_slide.notes_text_frame.text = (
        "NOTAS DO ORADOR (👤 INTEGRANTE 3):\n"
        "- Assumir a apresentação abordando a robustez contra falhas.\n"
        "- Responder às perguntas obrigatórias do edital: 'Como o robô identifica uma proposta inválida?' e 'Como o processo se recupera de uma falha?'.\n"
        "- Explicar o caso do Fornecedor D: detectado tanto por compliance web (Bloqueado) quanto por asserções numéricas (valores negativos).\n"
        "- Explicar que usamos o padrão Data Filtering Gateway para que propostas inválidas não distorçam o cálculo relativo do MCDA dos fornecedores válidos.\n"
        "- Demonstrar a resiliência multi-camada de rede que evita que o robô quebre se o servidor web cair."
    )

    # =========================================================================
    # SLIDE 8: TESTES AUTOMATIZADOS PYTEST (👤 Integrante 3)
    # =========================================================================
    slide8 = prs.slides.add_slide(blank_layout)
    aplicar_fundo_padrao(slide8, "Suíte de Testes Automatizados com Pytest", "7. Qualidade & Testes", "Integrante 3")

    c1 = adicionar_card(slide8, Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.9), "ESTRUTURA DA SUÍTE (51 TESTES)")
    tf = c1.text_frame
    p = tf.add_paragraph()
    p.text = "🧪 test_coleta.py (24 testes):\n" \
             "• Varredura de diretório e filtros de arquivos temporários (~$*, .*).\n" \
             "• Extração de fornecedores a partir de nomenclaturas.\n" \
             "• Playwright RPA com Chromium real e mocks de erro.\n" \
             "• Transição de abas no Playwright durante falhas de rede.\n\n" \
             "🧪 test_leitura.py (25 testes):\n" \
             "• Leitura de XLSX e CSV (delimitadores ';' e ',').\n" \
             "• Tolerância a encodings (utf-8, latin-1) e cabeçalhos com acentos.\n" \
             "• Leitura de critérios e template com fallback inteligente.\n\n" \
             "🧪 test_integration.py (2 testes):\n" \
             "• Teste de integração Coleta + Leitura + Main."
    p.font.size = Pt(11)
    p.font.color.rgb = COLOR_TEXT_MUTED

    c2 = adicionar_card(slide8, Inches(6.8), Inches(1.8), Inches(5.733), Inches(4.9), "COMANDO & RESULTADOS DE COBERTURA", bg_color=COLOR_BG_CARD_LIGHT)
    tf = c2.text_frame
    p = tf.add_paragraph()
    p.text = "💻 Execução dos Testes no Terminal:\n" \
             "   python -m pytest -v --cov=src\n\n" \
             "📊 Relatório de Execução:\n" \
             "• Total de Testes Executados: 51\n" \
             "• Testes Aprovados: 51 (100% Passed ✅)\n" \
             "• Falhas / Erros: 0\n" \
             "• Tempo de Execução: ~35 segundos\n\n" \
             "🎯 Cobertura de Código por Módulo:\n" \
             "• src/etapa1_coleta/:  100% de cobertura nos ramos\n" \
             "• src/etapa2_leitura/: 100% de cobertura nos ramos\n" \
             "• src/logger.py:       84% de cobertura\n\n" \
             "🛡️ Garantia de Não-Regressão em merges para develop e main."
    p.font.size = Pt(11)
    p.font.color.rgb = COLOR_TEXT_WHITE

    slide8.notes_slide.notes_text_frame.text = (
        "NOTAS DO ORADOR (👤 INTEGRANTE 3):\n"
        "- Apresentar a estratégia de garantia de qualidade através de testes automatizados.\n"
        "- Explicar que todos os requisitos do edital (fornecedor válido, fornecedor inválido, tratamento de dados negativos, cálculo de ranking) foram cobertos por testes unitários e de integração.\n"
        "- Mostrar que temos 51 testes implementados com 100% de sucesso e que a suíte roda automaticamente em cada push no GitHub Actions.\n"
        "- Passar para o próximo slide sobre o algoritmo matemático MCDA."
    )

    # =========================================================================
    # SLIDE 9: ALGORITMO MCDA PONDERADO (👤 Integrante 3)
    # =========================================================================
    slide9 = prs.slides.add_slide(blank_layout)
    aplicar_fundo_padrao(slide9, "Motor de Decisão Multicritério: Algoritmo MCDA", "8. Modelo Matemático", "Integrante 3")

    c1 = adicionar_card(slide9, Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.9), "PESOS DE NEGÓCIO & CRITÉRIOS")
    tf = c1.text_frame
    p = tf.add_paragraph()
    p.text = "📋 Matriz de Critérios (criterios_ranking.xlsx):\n\n" \
             "• Custo: Peso 40% (0.40) | Direção: Menor é Melhor 🔻\n" \
             "• Prazo (Lead Time): Peso 25% (0.25) | Direção: Menor é Melhor 🔻\n" \
             "• Capacidade: Peso 20% (0.20) | Direção: Maior é Melhor 🔺\n" \
             "• Qualidade: Peso 15% (0.15) | Direção: Maior é Melhor 🔺\n\n" \
             "⚖️ Ponderação:\n" \
             "   Soma dos Pesos = 0.40 + 0.25 + 0.20 + 0.15 = 1.00 (100%)\n\n" \
             "🔄 Carregamento Dinâmico:\n" \
             "• Lê da planilha oficial e aplica fallback seguro se o arquivo for excluído."
    p.font.size = Pt(11)
    p.font.color.rgb = COLOR_TEXT_MUTED

    c2 = adicionar_card(slide9, Inches(6.8), Inches(1.8), Inches(5.733), Inches(4.9), "FÓRMULAS DE NORMALIZAÇÃO RELATIVA", bg_color=COLOR_BG_CARD_LIGHT)
    tf = c2.text_frame
    p = tf.add_paragraph()
    p.text = "📐 Normalização Relativa ao Benchmark de Mercado:\n\n" \
             "1. Critérios de Menor Valor (Custo e Prazo):\n" \
             "   Score = (Menor Valor do Mercado / Valor da Proposta) * 100\n" \
             "   -> A melhor cotação recebe 100; as mais caras decrescem proporcionalmente.\n\n" \
             "2. Critérios de Maior Valor (Capacidade e Qualidade):\n" \
             "   Score = (Valor da Proposta / Maior Valor do Mercado) * 100\n" \
             "   -> A maior capacidade/qualidade recebe nota 100.\n\n" \
             "🎯 Equação do Score Global Ponderado:\n" \
             "   Nota Final = (Score_Custo * 0.40) + (Score_Prazo * 0.25) +\n" \
             "                (Score_Capac * 0.20) + (Score_Qual * 0.15)"
    p.font.size = Pt(11)
    p.font.color.rgb = COLOR_TEXT_WHITE

    slide9.notes_slide.notes_text_frame.text = (
        "NOTAS DO ORADOR (👤 INTEGRANTE 3):\n"
        "- Responder às perguntas: 'Como os pesos são aplicados?' e 'Como funciona a normalização dos valores?'.\n"
        "- Explicar o conceito de MCDA (Multi-Criteria Decision Analysis): evita decisões baseadas em um único atributo (ex: só preço baixo mas prazo péssimo).\n"
        "- Demonstrar a fórmula de normalização onde o melhor valor do mercado serve de benchmark (nota 100) e os demais são avaliados relativamente.\n"
        "- Passar a palavra para o Integrante 4, que explicará Docker, CI/CD, GHCR e Auditoria SOX."
    )

    # =========================================================================
    # SLIDE 10: DOCKER & DOCKER COMPOSE (👤 Integrante 4)
    # =========================================================================
    slide10 = prs.slides.add_slide(blank_layout)
    aplicar_fundo_padrao(slide10, "Containerização Docker & Docker Compose", "9. Infraestrutura & Containers", "Integrante 4")

    c1 = adicionar_card(slide10, Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.9), "DOCKERFILE & REQUISITOS EDITAL")
    tf = c1.text_frame
    p = tf.add_paragraph()
    p.text = "🐳 Imagem Base Python 3.12-slim:\n" \
             "• Otimizada e leve para ambiente de produção.\n\n" \
             "🕒 Timezone Obrigatório: TZ=America/Manaus\n" \
             "• Configurado no sistema operacional (tzdata) e nas variáveis de ambiente.\n" \
             "• Garante que os logs e carimbos de auditoria SOX sigam o horário de Manaus.\n\n" \
             "🌐 Instalação do Chromium Playwright:\n" \
             "• 'playwright install --with-deps chromium' embutido na construção da imagem.\n\n" \
             "📂 Volumes Persistentes:\n" \
             "• ./output -> /app/output (Planilha gerada)\n" \
             "• ./logs   -> /app/logs   (Auditoria JSON e logs)"
    p.font.size = Pt(11)
    p.font.color.rgb = COLOR_TEXT_MUTED

    c2 = adicionar_card(slide10, Inches(6.8), Inches(1.8), Inches(5.733), Inches(4.9), "ORQUESTRAÇÃO COM DOCKER COMPOSE", bg_color=COLOR_BG_CARD_LIGHT)
    tf = c2.text_frame
    p = tf.add_paragraph()
    p.text = "🚀 Execução Completa com 1 Único Comando:\n" \
             "   docker compose up --build\n\n" \
             "🛠️ Serviços Orquestrados:\n" \
             "1. Serviço 'web-panel':\n" \
             "   • Sobe servidor HTTP na porta 8000 simulando o portal corporativo da LG.\n" \
             "   • Healthcheck automático validando a disponibilidade do endpoint HTML.\n\n" \
             "2. Serviço 'robot':\n" \
             "   • Depende do 'web-panel' estar ativo ('depends_on').\n" \
             "   • Executa a esteira ponta a ponta (ETL, MCDA e exportação).\n\n" \
             "🔒 Isolamento total sem necessidade de instalar dependências na máquina host."
    p.font.size = Pt(11)
    p.font.color.rgb = COLOR_TEXT_WHITE

    slide10.notes_slide.notes_text_frame.text = (
        "NOTAS DO ORADOR (👤 INTEGRANTE 4):\n"
        "- Assumir a apresentação abordando a infraestrutura de DevOps e Containerização.\n"
        "- Explicar como atendemos a todos os requisitos do edital: Dockerfile, Docker Compose, Volumes persistentes, Variáveis de ambiente e Timezone America/Manaus.\n"
        "- Destacar a arquitetura multi-serviço no Docker Compose: o serviço web-panel disponibiliza o portal simulado e o serviço robot executa a esteira consumindo o portal e gerando os artefatos nos volumes mapeados."
    )

    # =========================================================================
    # SLIDE 11: CI/CD GITHUB ACTIONS & GHCR (👤 Integrante 4)
    # =========================================================================
    slide11 = prs.slides.add_slide(blank_layout)
    aplicar_fundo_padrao(slide11, "Pipeline de CI/CD (GitHub Actions) & GHCR", "10. DevOps & Integração Contínua", "Integrante 4")

    c1 = adicionar_card(slide11, Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.9), "PIPELINE GITHUB ACTIONS (.github/workflows/ci.yml)")
    tf = c1.text_frame
    p = tf.add_paragraph()
    p.text = "🔄 Gatilhos Automáticos:\n" \
             "• Push e Pull Request nas branches 'main', 'develop' e 'feature/**'.\n\n" \
             "⚙️ Job 1: Lint, Testes e Cobertura:\n" \
             "1. Checkout do código e Setup Python 3.12 com cache pip.\n" \
             "2. Instalação de dependências e Chromium do Playwright.\n" \
             "3. Verificação de Linting estrito com Flake8 (0 erros).\n" \
             "4. Execução dos 51 testes Pytest com relatório de cobertura XML/Term.\n\n" \
             "⚙️ Job 2: Build e Validação Docker:\n" \
             "• Executa após a aprovação dos testes ('needs: lint-and-test').\n" \
             "• Valida a integridade da imagem Docker."
    p.font.size = Pt(11)
    p.font.color.rgb = COLOR_TEXT_MUTED

    c2 = adicionar_card(slide11, Inches(6.8), Inches(1.8), Inches(5.733), Inches(4.9), "PUBLICAÇÃO & EXECUÇÃO VIA GHCR", bg_color=COLOR_BG_CARD_LIGHT)
    tf = c2.text_frame
    p = tf.add_paragraph()
    p.text = "📦 GitHub Container Registry (GHCR):\n" \
             "• Imagem homologada publicada em ghcr.io.\n" \
             "• Não utilizamos Docker Hub, atendendo estritamente ao edital da disciplina.\n\n" \
             "💻 Comandos de Demonstração em Produção:\n" \
             "1. Baixar a imagem publicada:\n" \
             "   docker pull ghcr.io/sannyer3232/hyperautomation-av3-equipe1:latest\n\n" \
             "2. Executar o container com mapeamento de volumes:\n" \
             "   docker run --rm -v $(pwd)/output:/app/output \\\n" \
             "              -v $(pwd)/logs:/app/logs \\\n" \
             "              -e TZ=America/Manaus \\\n" \
             "              ghcr.io/sannyer3232/hyperautomation-av3-equipe1:latest\n\n" \
             "✅ Execução autônoma e geração imediata dos relatórios."
    p.font.size = Pt(11)
    p.font.color.rgb = COLOR_TEXT_WHITE

    slide11.notes_slide.notes_text_frame.text = (
        "NOTAS DO ORADOR (👤 INTEGRANTE 4):\n"
        "- Explicar o ciclo completo de CI/CD: Código -> Push -> Flake8 -> Pytest -> Docker Build -> Publicação GHCR.\n"
        "- Ressaltar que a esteira valida automaticamente a qualidade do código antes de permitir qualquer merge para o develop ou main.\n"
        "- Demonstrar como a aplicação pode ser puxada diretamente do GHCR via 'docker pull' e executada via 'docker run' em qualquer servidor da LG no mundo."
    )

    # =========================================================================
    # SLIDE 12: LOGS E AUDITORIA SOX (👤 Integrante 4)
    # =========================================================================
    slide12 = prs.slides.add_slide(blank_layout)
    aplicar_fundo_padrao(slide12, "Rastreabilidade Digital, Logs e Auditoria SOX", "11. Governança & Compliance", "Integrante 4")

    c1 = adicionar_card(slide12, Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.9), "TRILHA DE AUDITORIA (logs/auditoria.json)")
    tf = c1.text_frame
    p = tf.add_paragraph()
    p.text = "🏛️ Snapshot Estruturado para Compliance SOX:\n\n" \
             "• Data/Hora de Início e Fim (ISO 8601 com TZ Manaus).\n" \
             "• Total de propostas recebidas (4) e arquivos lidos.\n" \
             "• Lista de propostas aprovadas com dados comerciais.\n" \
             "• Lista de propostas rejeitadas com justificativa técnica (Fornecedor D).\n" \
             "• Memória de cálculo detalhada com scores parciais por critério (Custo, Prazo, Capacidade, Qualidade).\n" \
             "• Ranking final homologado com classificação oficial.\n" \
             "• Registro de erros e exceções capturadas."
    p.font.size = Pt(11)
    p.font.color.rgb = COLOR_TEXT_MUTED

    c2 = adicionar_card(slide12, Inches(6.8), Inches(1.8), Inches(5.733), Inches(4.9), "LOGS OPERACIONAIS & ARTEFATO FINAL", bg_color=COLOR_BG_CARD_LIGHT)
    tf = c2.text_frame
    p = tf.add_paragraph()
    p.text = "📋 Log Operacional Contínuo (logs/execucao.log):\n" \
             "• Formatação padronizada: %(asctime)s [%(levelname)s] [%(name)s] %(message)s\n" \
             "• 4 Níveis de Severidade Utilizados:\n" \
             "  - [INFO]: Etapas normais, propostas lidas e ranking gerado.\n" \
             "  - [WARNING]: Fallback de rede ativado, proposta D rejeitada.\n" \
             "  - [ERROR]: Falha técnica localizada na leitura de arquivos.\n" \
             "  - [CRITICAL]: Erro fatal não tratado que impeça a execução.\n\n" \
             "📊 Planilha Final Homologada (output/ranking_final.xlsx):\n" \
             "• Preenchimento fiel do template modelo_ranking.xlsx.\n" \
             "• Posições: 1º (Fornec. B), 2º (Fornec. A), 3º (Fornec. C), Desclassificado (Fornec. D)."
    p.font.size = Pt(11)
    p.font.color.rgb = COLOR_TEXT_WHITE

    slide12.notes_slide.notes_text_frame.text = (
        "NOTAS DO ORADOR (👤 INTEGRANTE 4):\n"
        "- Responder à pergunta: 'Como o resultado é auditado?'.\n"
        "- Explicar a diferença entre o log operacional (execucao.log) para diagnóstico de TI e a trilha digital de auditoria (auditoria.json) para governança corporativa e compliance SOX.\n"
        "- Mostrar que qualquer auditor externo pode abrir o JSON de auditoria e auditar exatamente quem cotou, qual score cada critério recebeu, por que o vencedor ganhou e por que o Fornecedor D foi reprovado."
    )

    # =========================================================================
    # SLIDE 13: CONCLUSÃO E VALOR ENTREGUE (👤 Integrante 4)
    # =========================================================================
    slide13 = prs.slides.add_slide(blank_layout)
    aplicar_fundo_padrao(slide13, "Conclusão e Entrega de Valor para o Negócio", "12. Fechamento & ROI", "Integrante 4")

    c1 = adicionar_card(slide13, Inches(0.8), Inches(1.8), Inches(3.6), Inches(4.9), "O QUE FOI DESENVOLVIDO?")
    tf = c1.text_frame
    p = tf.add_paragraph()
    p.text = "• Solução completa de Hyperautomation.\n" \
             "• RPA Playwright + Motor MCDA + Pytest + Docker + CI/CD + GHCR.\n" \
             "• 100% dos requisitos do edital da LG atendidos."
    p.font.size = Pt(11)
    p.font.color.rgb = COLOR_TEXT_MUTED

    c2 = adicionar_card(slide13, Inches(4.8), Inches(1.8), Inches(3.7), Inches(4.9), "QUAL PROBLEMA RESOLVEU?")
    tf = c2.text_frame
    p = tf.add_paragraph()
    p.text = "• Eliminação de 100% da digitação manual de propostas.\n" \
             "• Barramento automático de fornecedores irregulares.\n" \
             "• Redução do ciclo de compras de 8 horas para 25 segundos."
    p.font.size = Pt(11)
    p.font.color.rgb = COLOR_TEXT_MUTED

    c3 = adicionar_card(slide13, Inches(8.9), Inches(1.8), Inches(3.633), Inches(4.9), "QUAL VALOR ENTREGA?", bg_color=COLOR_BG_CARD_LIGHT)
    tf = c3.text_frame
    p = tf.add_paragraph()
    p.text = "• Decisões 100% objetivas e matemáticas.\n" \
             "• Total conformidade com auditorias SOX da LG.\n" \
             "• Escalabilidade para suportar mais de 350 cotações/mês na fábrica de Manaus."
    p.font.size = Pt(11)
    p.font.color.rgb = COLOR_TEXT_WHITE

    slide13.notes_slide.notes_text_frame.text = (
        "NOTAS DO ORADOR (👤 INTEGRANTE 4):\n"
        "- Fazer a síntese executiva do projeto respondendo às 5 perguntas finais da apresentação:\n"
        "  1. O que foi desenvolvido?\n"
        "  2. Qual problema foi resolvido?\n"
        "  3. Como a solução funciona?\n"
        "  4. Qual resultado foi obtido?\n"
        "  5. Qual valor a solução entrega para o processo da LG?\n"
        "- Concluir a parte teórica e anunciar o início da demonstração prática ao vivo."
    )

    # =========================================================================
    # SLIDE 14: ROTEIRO DA DEMONSTRAÇÃO PRÁTICA (👥 4 Integrantes)
    # =========================================================================
    slide14 = prs.slides.add_slide(blank_layout)
    aplicar_fundo_padrao(slide14, "Roteiro da Demonstração Prática (7 Etapas)", "Demonstração ao Vivo", "4 Integrantes")

    c1 = adicionar_card(slide14, Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.9), "ETAPAS 1 A 4 DA DEMONSTRAÇÃO")
    tf = c1.text_frame
    p = tf.add_paragraph()
    p.text = "👤 Integrante 1 — ETAPA 1: Código & Repositório\n" \
             "• Mostrar organização de pastas (src, tests, resources, docs, .github).\n" \
             "• Explicar arquivos de configuração (.env, config.py, logger.py).\n\n" \
             "👤 Integrante 2 — ETAPA 2: Execução da Pipeline\n" \
             "• Executar 'python src/main.py' no terminal.\n" \
             "• Mostrar no log a coleta Playwright e a leitura multiformato.\n\n" \
             "👤 Integrante 3 — ETAPA 3: Testes Automatizados\n" \
             "• Executar 'python -m pytest -v --tb=short'.\n" \
             "• Demonstrar os 51 testes passando com 100% de sucesso.\n\n" \
             "👤 Integrante 4 — ETAPA 4: Workflow no GitHub Actions\n" \
             "• Abrir a aba Actions no GitHub e mostrar a esteira verde (Lint + Test + Docker)."
    p.font.size = Pt(11)
    p.font.color.rgb = COLOR_TEXT_MUTED

    c2 = adicionar_card(slide14, Inches(6.8), Inches(1.8), Inches(5.733), Inches(4.9), "ETAPAS 5 A 7 DA DEMONSTRAÇÃO", bg_color=COLOR_BG_CARD_LIGHT)
    tf = c2.text_frame
    p = tf.add_paragraph()
    p.text = "👤 Integrante 4 — ETAPA 5: Imagem no GHCR\n" \
             "• Mostrar o pacote publicado no GitHub Container Registry (ghcr.io).\n\n" \
             "👤 Integrante 4 — ETAPA 6: Execução via Docker / Compose\n" \
             "• Demonstrar 'docker compose up' ou 'docker run' da imagem do GHCR.\n" \
             "• Mostrar a execução em container com timezone America/Manaus.\n\n" \
             "👥 Todos — ETAPA 7: Resultados, Logs & Auditoria\n" \
             "• Abrir e exibir o arquivo 'output/ranking_final.xlsx' preenchido.\n" \
             "• Abrir o 'logs/auditoria.json' mostrando a memória de cálculo e desclassificação do Fornecedor D.\n" \
             "• Abrir o 'logs/execucao.log' demonstrando severidades INFO e WARNING."
    p.font.size = Pt(11)
    p.font.color.rgb = COLOR_TEXT_WHITE

    slide14.notes_slide.notes_text_frame.text = (
        "NOTAS DO ORADOR (👥 TODOS OS 4 INTEGRANTES):\n"
        "- Seguir exatamente este roteiro cronometrado durante a demonstração ao vivo para a banca:\n"
        "  1. Integrante 1 mostra o código e a estrutura do repositório no VS Code / GitHub.\n"
        "  2. Integrante 2 executa o 'python src/main.py' no terminal mostrando a execução do robô.\n"
        "  3. Integrante 3 executa 'python -m pytest -v' demonstrando os testes automatizados.\n"
        "  4. Integrante 4 mostra o GitHub Actions, o pacote publicado no GHCR e executa no Docker.\n"
        "  5. Equipe mostra a planilha final gerada e o JSON de auditoria SOX."
    )

    # =========================================================================
    # SLIDE 15: GUIA DE DEFESA TÉCNICA INDIVIDUAL (👥 4 Integrantes)
    # =========================================================================
    slide15 = prs.slides.add_slide(blank_layout)
    aplicar_fundo_padrao(slide15, "Defesa Técnica Individual: Perguntas da Banca & Dúvidas", "Perguntas & Respostas", "4 Integrantes")

    c1 = adicionar_card(slide15, Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.9), "PERGUNTAS FREQUENTES DA BANCA")
    tf = c1.text_frame
    p = tf.add_paragraph()
    p.text = "1. 'Como identifica uma proposta inválida?'\n" \
             "-> Dupla validação: Status Web (Bloqueado) e dados negativos (<=0).\n\n" \
             "2. 'Como os pesos e o ranking são calculados?'\n" \
             "-> Carga de criterios_ranking.xlsx (40/25/20/15) + normalização MCDA.\n\n" \
             "3. 'Como evita que o Fornecedor D contamine o ranking?'\n" \
             "-> Data Filtering Gateway: segregado antes do cálculo do benchmark.\n\n" \
             "4. 'Como o robô se recupera de falhas?'\n" \
             "-> Resiliência multi-camada: Playwright HTTP -> file:// -> Requests -> Disco.\n\n" \
             "5. 'Como o resultado é auditado?'\n" \
             "-> Trilha estruturada em logs/auditoria.json para conformidade SOX."
    p.font.size = Pt(10)
    p.font.color.rgb = COLOR_TEXT_MUTED

    c2 = adicionar_card(slide15, Inches(6.8), Inches(1.8), Inches(5.733), Inches(4.9), "OBRIGADO! DÚVIDAS DA BANCA?", bg_color=COLOR_BG_CARD_LIGHT)
    tf = c2.text_frame
    p = tf.add_paragraph()
    p.text = "LG ELECTRONICS DO BRASIL\n" \
             "Automação Inteligente do Processo de Seleção de Fornecedores\n\n" \
             "👨‍💻 Equipe 01 — Turma T02 (2026)\n" \
             "• Sannyer Carvalho & Integrantes da Equipe\n\n" \
             "🎓 Disciplina: Técnicas de Hyperautomation\n" \
             "👨‍🏫 Professor Responsável: Prof. Moisés Levy\n\n" \
             "🔗 Repositório GitHub: github.com/Sannyer3232/hyperautomation-av3-equipe1\n" \
             "📦 Docker Image (GHCR): ghcr.io/sannyer3232/hyperautomation-av3-equipe1:latest\n\n" \
             "Estamos abertos para as perguntas da banca avaliadora!"
    p.font.size = Pt(11)
    p.font.color.rgb = COLOR_TEXT_WHITE

    slide15.notes_slide.notes_text_frame.text = (
        "NOTAS DO ORADOR (👥 TODOS OS 4 INTEGRANTES):\n"
        "- Agradecer a atenção do professor Moisés Levy e da turma.\n"
        "- Informar que a equipe está pronta para responder a qualquer pergunta técnica individual ou coletiva sobre Python, Playwright, MCDA, Pytest, Docker, CI/CD, GHCR ou Auditoria SOX."
    )

    prs.save(caminho_saida)
    print(f"Apresentação gerada com sucesso em: {caminho_saida}")


if __name__ == "__main__":
    output_path = Path("docs") / "Apresentacao_Defesa_Hyperautomation_LG_Equipe01.pptx"
    criar_apresentacao_completa(str(output_path))
