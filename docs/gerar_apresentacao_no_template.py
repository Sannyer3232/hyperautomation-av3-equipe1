"""
Script para geração da Apresentação de Defesa Técnica da Avaliação 3
utilizando EXATAMENTE o template e a identidade visual de:
'docs/Apresentação DM Sannyer Cardoso.pdf (2) (1).pptx'.

Garante 15 slides completos, com divisão para 4 apresentadores,
design fiel ao template original da AX Academy / Hyperautomation,
notas de orador completas e formatação profissional.
"""

import copy
from pathlib import Path
from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Dimensões exatas do template original (EMU)
SLIDE_WIDTH = 5765800
SLIDE_HEIGHT = 3238500

# Cores em harmonia com o template original e a identidade LG
COLOR_PRIMARY = RGBColor(165, 0, 52)       # LG Ruby Red
COLOR_CYAN = RGBColor(0, 180, 216)          # Electric Cyan
COLOR_GOLD = RGBColor(255, 209, 102)        # Gold / Destaque
COLOR_GREEN = RGBColor(6, 214, 160)         # Emerald Green
COLOR_DARK_TEXT = RGBColor(30, 41, 59)      # Slate Dark
COLOR_MUTED_TEXT = RGBColor(71, 85, 105)    # Slate Muted
COLOR_CARD_BG = RGBColor(241, 245, 249)     # Slate Light Background
COLOR_CARD_BORDER = RGBColor(203, 213, 225) # Border Slate
COLOR_WHITE = RGBColor(255, 255, 255)


def clonar_slide_template(prs, source_slide_idx):
    """
    Clona um slide do template preservando elementos gráficos,
    imagens de fundo, linhas e estrutura de cabeçalho/rodapé.
    """
    source_slide = prs.slides[source_slide_idx]
    slide_layout = source_slide.slide_layout
    new_slide = prs.slides.add_slide(slide_layout)
    
    # Copia todas as formas do slide original
    for shape in source_slide.shapes:
        new_slide.shapes._spTree.append(copy.deepcopy(shape._element))
        
    return new_slide


def limpar_formas_conteudo_antigo(slide):
    """
    Remove apenas os blocos de texto de conteúdo antigo do absenteísmo,
    preservando as formas de cabeçalho, linhas, logotipos e rodapé.
    """
    formas_para_remover = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            texto = shape.text_frame.text.lower()
            if any(termo in texto for termo in [
                "absenteeism", "social drinker", "demográficas", "anomalias",
                "imputação", "categóricos", "tree regressor", "viabilidade promissora",
                "features", "target", "refinamento"
            ]):
                formas_para_remover.append(shape)
        elif shape.has_table:
            formas_para_remover.append(shape)

    for shape in formas_para_remover:
        sp = shape._element
        sp.getparent().remove(sp)


def ajustar_cabecalho_e_rodape(slide, titulo_texto, numero_slide_str, integrante_str=""):
    """
    Atualiza o título do slide e os placeholders do rodapé original.
    """
    for shape in slide.shapes:
        if shape.has_text_frame:
            tf = shape.text_frame
            texto = tf.text.strip()
            
            # Atualiza título do cabeçalho
            if shape.top < Emu(400000) and shape.left > Emu(500000):
                tf.clear()
                p = tf.paragraphs[0]
                p.text = titulo_texto
                p.font.size = Pt(13)
                p.font.bold = True
                p.font.color.rgb = COLOR_PRIMARY
                
            # Atualiza número da página no rodapé
            elif "‹#›" in texto or "#" in texto or "/9" in texto:
                tf.clear()
                p = tf.paragraphs[0]
                p.text = numero_slide_str
                p.font.size = Pt(7)
                p.font.color.rgb = COLOR_MUTED_TEXT

            # Atualiza autor / equipe no rodapé
            elif "sannyer" in texto.lower() and shape.top > Emu(2800000):
                tf.clear()
                p = tf.paragraphs[0]
                p.text = "Equipe 01 — Automação Seleção de Fornecedores LG"
                p.font.size = Pt(7)
                p.font.color.rgb = COLOR_MUTED_TEXT

            # Atualiza disciplina / instituição no rodapé
            elif "modelo de apresentação" in texto.lower():
                tf.clear()
                p = tf.paragraphs[0]
                p.text = "Técnicas de Hyperautomation • Prof. Moisés Levy (T02)"
                p.font.size = Pt(7)
                p.font.color.rgb = COLOR_MUTED_TEXT

    # Adiciona Badge do Integrante Responsável no canto superior direito
    if integrante_str:
        badge_w, badge_h = Emu(1600000), Emu(190000)
        badge_x, badge_y = Emu(SLIDE_WIDTH - 1800000), Emu(170000)
        badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, badge_x, badge_y, badge_w, badge_h)
        badge.fill.solid()
        badge.fill.fore_color.rgb = COLOR_PRIMARY
        badge.line.fill.background()
        tf_b = badge.text_frame
        tf_b.vertical_anchor = MSO_ANCHOR.MIDDLE
        p_b = tf_b.paragraphs[0]
        p_b.text = f"👤 {integrante_str}"
        p_b.font.size = Pt(7.5)
        p_b.font.bold = True
        p_b.font.color.rgb = COLOR_WHITE
        p_b.alignment = PP_ALIGN.CENTER


def adicionar_container_card(slide, x, y, w, h, titulo="", bg_color=COLOR_CARD_BG, border_color=COLOR_CARD_BORDER):
    """Cria um card com fundo suave e borda estruturada."""
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    card.fill.solid()
    card.fill.fore_color.rgb = bg_color
    card.line.color.rgb = border_color
    card.line.width = Pt(1)

    tf = card.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(100000)
    tf.margin_top = Emu(80000)
    tf.margin_right = Emu(100000)
    tf.margin_bottom = Emu(80000)

    if titulo:
        p_tit = tf.paragraphs[0]
        p_tit.text = titulo
        p_tit.font.size = Pt(9.5)
        p_tit.font.bold = True
        p_tit.font.color.rgb = COLOR_PRIMARY
        p_tit.space_after = Pt(4)

    return card


def gerar_apresentacao_oficial():
    template_path = Path("docs") / "Apresentação DM Sannyer Cardoso.pdf (2) (1).pptx"
    prs = Presentation(str(template_path))

    # Guardamos os slides de referência do template original
    slide_conteudo_ref_idx = 7  # Slide 8 do template (layout limpo com linha e rodapé)
    slide_fim_ref_idx = 8       # Slide 9 do template (Obrigado / Dúvidas)

    # =========================================================================
    # SLIDE 1: CAPA DO PROJETO (👤 Integrante 1)
    # =========================================================================
    s1 = prs.slides[0]
    for shape in s1.shapes:
        if shape.has_text_frame:
            texto = shape.text_frame.text
            if "Modelagem Preditiva" in texto or "DATA SCIENCE" in texto:
                shape.text_frame.clear()
                p1 = shape.text_frame.paragraphs[0]
                p1.text = "Automação Inteligente do Processo de\nSeleção de Fornecedores Industriais"
                p1.font.size = Pt(19)
                p1.font.bold = True
                p1.font.color.rgb = COLOR_WHITE
                
                p2 = shape.text_frame.add_paragraph()
                p2.text = "HYPERAUTOMATION, RPA (PLAYWRIGHT), MOTOR MCDA, DOCKER & CI/CD — LG ELECTRONICS"
                p2.font.size = Pt(9.0)
                p2.font.color.rgb = COLOR_CYAN
                p2.space_before = Pt(8)
            elif "Sannyer Cardoso" in texto or "Professora:" in texto:
                shape.text_frame.clear()
                p = shape.text_frame.paragraphs[0]
                p.text = "Equipe 01 — Turma T02 (Agosto/2026)\n" \
                         "Integrantes: Sannyer Carvalho, Integrante 2, Integrante 3, Integrante 4\n" \
                         "Professor Responsável: Prof. Moisés Levy"
                p.font.size = Pt(8.5)
                p.font.color.rgb = COLOR_WHITE
            elif "AX Academy" in texto:
                shape.text_frame.clear()
                p = shape.text_frame.paragraphs[0]
                p.text = "AX Academy • Digital Transformation & Hyperautomation"
                p.font.size = Pt(8.5)
                p.font.color.rgb = COLOR_CYAN

    s1.notes_slide.notes_text_frame.text = (
        "NOTAS DO ORADOR (👤 INTEGRANTE 1):\n"
        "- Cumprimentar o professor Moisés Levy e a banca avaliadora.\n"
        "- Apresentar a Equipe 01 e introduzir o projeto: Automação Inteligente de Seleção de Fornecedores para a LG Electronics.\n"
        "- Destacar que a apresentação segue a narrativa oficial: Problema -> Solução -> Automação -> Testes -> Docker -> CI/CD -> GHCR -> Execução -> Resultado.\n"
        "- Explicar que a defesa está dividida equilibradamente entre 4 integrantes."
    )

    # =========================================================================
    # SLIDES 2 A 14: CONTEÚDO TÉCNICO E ARQUITETURA
    # =========================================================================
    dados_slides = [
        # SLIDE 2: O Problema AS-IS (👤 Integrante 1)
        {
            "titulo": "O Problema de Negócio: Processo Manual AS-IS na LG",
            "integrante": "Integrante 1 (Líder / Negócio)",
            "cards": [
                ("1. FLUXO OPERACIONAL MANUAL (AS-IS)",
                 "• Disparo Manual de RFQs: Compradores enviam cotações por e-mail para dezenas de fornecedores.\n"
                 "• Recepção Despadronizada: Propostas chegam em formatos dispersos (.xlsx, .csv, corpos de e-mail).\n"
                 "• Transcrição no Excel: Comprador digita preços, prazos e tributos manualmente em planilhas locais.\n"
                 "• Checagem Fragmentada: Consulta manual de status em portais externos.\n"
                 "• Cálculo Manual de TCO: Fórmulas manuais suscetíveis a erros antes do parecer da gerência.",
                 Emu(350000), Emu(550000), Emu(2450000), Emu(2150000)),
                ("2. IMPACTOS & RISCOS OPERACIONAIS",
                 "🚨 Erros Humanos e Fórmulas Inconsistentes:\n"
                 "• Até 12% de divergências em transcrição de valores e prazos.\n"
                 "• Risco de homologação de fornecedores bloqueados ou com dados negativos.\n\n"
                 "⏱️ Alto Lead Time de Suprimentos:\n"
                 "• Ciclo de compras levava de 5 a 10 dias úteis (~8h de esforço manual).\n\n"
                 "🔍 Baixa Rastreabilidade e Falta de Auditoria SOX:\n"
                 "• Inexistência de logs estruturados e histórico transparente das decisões.",
                 Emu(2900000), Emu(550000), Emu(2500000), Emu(2150000))
            ],
            "notes": "NOTAS DO ORADOR (👤 INTEGRANTE 1):\n- Explicar o cenário AS-IS da fábrica da LG em Manaus.\n- Destacar os gargalos de transcrição manual de planilhas Excel e arquivos CSV.\n- Enfatizar os três grandes riscos: erro humano, lead time de até 10 dias e vulnerabilidade em auditorias SOX."
        },

        # SLIDE 3: A Solução TO-BE (👤 Integrante 1)
        {
            "titulo": "A Solução Proposta: Hyperautomation de Ponta a Ponta",
            "integrante": "Integrante 1 (Líder / Negócio)",
            "cards": [
                ("1. INGESTÃO & RPA",
                 "• Varredura autônoma de diretórios (.xlsx e .csv).\n"
                 "• Robô Playwright acessa portal web simulado.\n"
                 "• Extração de status cadastrais (Ativo/Bloqueado).\n"
                 "• Resiliência com fallback multi-camada.",
                 Emu(350000), Emu(550000), Emu(1600000), Emu(2150000)),
                ("2. VALIDAÇÃO & DECISÃO",
                 "• Gateway de validação cadastral e numérica.\n"
                 "• Rejeição instantânea de dados negativos.\n"
                 "• Algoritmo MCDA ponderado com normalização relativa de mercado.\n"
                 "• Pesos dinâmicos oficiais.",
                 Emu(2030000), Emu(550000), Emu(1650000), Emu(2150000)),
                ("3. GOVERNANÇA & AUDITORIA",
                 "• Preenchimento do template oficial da LG.\n"
                 "• Geração do ranking_final.xlsx.\n"
                 "• Trilha digital SOX em logs/auditoria.json.\n"
                 "• Logs contínuos com 4 níveis de severidade.",
                 Emu(3760000), Emu(550000), Emu(1650000), Emu(2150000))
            ],
            "notes": "NOTAS DO ORADOR (👤 INTEGRANTE 1):\n- Apresentar a visão TO-BE dividida em Ingestão/RPA, Validação/Decisão e Governança/Auditoria.\n- Enfatizar a redução de 95% do tempo de processamento e conformidade de 100% com auditorias.\n- Passar a palavra para o Integrante 2."
        },

        # SLIDE 4: Fluxo do Pipeline (👤 Integrante 2)
        {
            "titulo": "Fluxo da Solução: Pipeline Modular em 6 Etapas",
            "integrante": "Integrante 2 (Automação / Dados)",
            "cards": [
                ("ETAPA 1: COLETA", "• Varredura .xlsx/.csv\n• Playwright RPA\n• Portal Web HTTP\n• Fallback local file://", Emu(350000), Emu(550000), Emu(820000), Emu(2150000)),
                ("ETAPA 2: LEITURA", "• Leitura Pandas\n• Tolerância delimitadores\n• Normalização colunas\n• Carga de critérios", Emu(1220000), Emu(550000), Emu(820000), Emu(2150000)),
                ("ETAPA 3: VALIDAÇÃO", "• Checagem portal web\n• Barramento negativo\n• Rejeição Fornec. D\n• Data Filtering", Emu(2090000), Emu(550000), Emu(820000), Emu(2150000)),
                ("ETAPA 4: CONSOLIDAÇÃO", "• Unificação DataFrames\n• Tipagem estrita\n• Isolamento rejeitadas\n• Preparação MCDA", Emu(2960000), Emu(550000), Emu(820000), Emu(2150000)),
                ("ETAPA 5: RANKING", "• Scoring relativo\n• Pesos ponderados\n• Menor custo/prazo\n• Maior capac./qualid.", Emu(3830000), Emu(550000), Emu(820000), Emu(2150000)),
                ("ETAPA 6: RESULTADO", "• Template oficial\n• ranking_final.xlsx\n• Auditoria SOX JSON\n• Logs detalhados", Emu(4700000), Emu(550000), Emu(710000), Emu(2150000))
            ],
            "notes": "NOTAS DO ORADOR (👤 INTEGRANTE 2):\n- Assumir a apresentação explicando a arquitetura em 6 etapas da pipeline.\n- Explicar como cada módulo possui responsabilidade bem definida e desacoplada, facilitando testes e manutenção."
        },

        # SLIDE 5: Stack Tecnológico (👤 Integrante 2)
        {
            "titulo": "Stack Tecnológico e Integração dos Componentes",
            "integrante": "Integrante 2 (Automação / Dados)",
            "cards": [
                ("LINGUAGEM & ROBÓTICA",
                 "🐍 Python 3.12 (Core Engine):\n"
                 "• Orquestrador modular com tipagem estrita.\n\n"
                 "🌐 Playwright (Chromium RPA Headless):\n"
                 "• Automação de navegador ultrarrápida.\n"
                 "• Extração de tabelas e status no portal web.\n\n"
                 "📊 Pandas & OpenPyXL:\n"
                 "• Processamento de dados e preenchimento de templates Excel.",
                 Emu(350000), Emu(550000), Emu(1600000), Emu(2150000)),
                ("QUALIDADE & TESTES",
                 "🧪 Pytest (v8.0+):\n"
                 "• Suíte com 51 testes automatizados (unitários e integração).\n\n"
                 "📈 Pytest-Cov:\n"
                 "• Relatórios de cobertura de código nos módulos.\n\n"
                 "🧹 Flake8:\n"
                 "• Padronização estrita PEP8 com 0 erros de linting.",
                 Emu(2030000), Emu(550000), Emu(1650000), Emu(2150000)),
                ("INFRAESTRUTURA & DEVOPS",
                 "🐳 Docker & Docker Compose:\n"
                 "• Containers isolados (Robô + Web Server).\n"
                 "• Timezone oficial: America/Manaus.\n\n"
                 "⚙️ GitHub Actions & GHCR:\n"
                 "• Esteira de CI/CD completa.\n"
                 "• Publicação no GitHub Container Registry (ghcr.io).",
                 Emu(3760000), Emu(550000), Emu(1650000), Emu(2150000))
            ],
            "notes": "NOTAS DO ORADOR (👤 INTEGRANTE 2):\n- Justificar as escolhas técnicas: Playwright moderno para RPA web headless, Pandas para operações vetoriais de MCDA, Pytest para testes e Docker Compose com timezone America/Manaus."
        },

        # SLIDE 6: Cenário Normal / Caminho Feliz (👤 Integrante 2)
        {
            "titulo": "Cenário de Execução Normal: Fornecedores A, B e C",
            "integrante": "Integrante 2 (Automação / Dados)",
            "cards": [
                ("1. ENTRADA & PROCESSAMENTO NORMAL",
                 "• Fornecedor A (.xlsx): Custo R$ 100 | Prazo: 5d | Cap: 500 | Qual: 95\n"
                 "• Fornecedor B (.csv):  Custo R$ 90  | Prazo: 8d | Cap: 700 | Qual: 90\n"
                 "• Fornecedor C (.xlsx): Custo R$ 110 | Prazo: 4d | Cap: 400 | Qual: 98\n\n"
                 "🌐 Consulta Cadastral Playwright:\n"
                 "• Fornecedores A, B e C validados com status 'Ativo' no portal web.\n\n"
                 "✅ Validação Numérica Aprovada:\n"
                 "• Valores positivos e dentro dos limites operacionais da LG.\n\n"
                 "📊 Consolidação:\n"
                 "• 100% dos dados válidos encaminhados para o cálculo MCDA.",
                 Emu(350000), Emu(550000), Emu(2450000), Emu(2150000)),
                ("2. RESULTADO & HOMOLOGAÇÃO DO CAMINHO FELIZ",
                 "🏆 Homologação Automática e Ranking Gerado:\n\n"
                 "• 1º Lugar: Fornecedor B (Score Final: 92.98)\n"
                 "  -> Menor custo de mercado (R$ 90) e maior capacidade produtiva (700 un).\n\n"
                 "• 2º Lugar: Fornecedor A (Score Final: 89.28)\n"
                 "  -> Proposta equilibrada entre prazo (5d) e qualidade (95%).\n\n"
                 "• 3º Lugar: Fornecedor C (Score Final: 88.08)\n"
                 "  -> Melhor prazo (4d) e maior qualidade (98%), porém custo mais alto (R$ 110).\n\n"
                 "📁 Artefato Gerado: output/ranking_final.xlsx.",
                 Emu(2900000), Emu(550000), Emu(2500000), Emu(2150000))
            ],
            "notes": "NOTAS DO ORADOR (👤 INTEGRANTE 2):\n- Demonstrar o caminho feliz da esteira.\n- Mostrar que os 3 fornecedores válidos foram processados e o Fornecedor B venceu pelo peso preponderante de Custo (40%).\n- Passar a palavra para o Integrante 3."
        },

        # SLIDE 7: Tratamento de Exceções e Resiliência (👤 Integrante 3)
        {
            "titulo": "Tratamento de Exceções, Resiliência e Barramentos",
            "integrante": "Integrante 3 (Qualidade / MCDA)",
            "cards": [
                ("CASO 1: DESCLASSIFICAÇÃO FORNECEDOR D",
                 "🚨 Dupla Barreira de Não-Conformidade:\n"
                 "1. Camada Cadastral Web: Status 'Bloqueado' detectado no portal HTML.\n"
                 "2. Camada Numérica: Custo negativo (-50), Prazo (-2) e Capacidade (-100).\n\n"
                 "🛡️ Padrão Data Filtering Gateway:\n"
                 "• A proposta D é imediatamente segregada em 'propostas_rejeitadas'.\n"
                 "• Não contamina o cálculo relativo do MCDA dos concorrentes.\n"
                 "• Figura no ranking oficial como 'Desclassificado' (Nota 0.00) com justificativa auditada.",
                 Emu(350000), Emu(550000), Emu(2450000), Emu(2150000)),
                ("CASO 2: RESILIÊNCIA DE REDE & FALLBACK",
                 "🔄 Arquitetura Multi-Camada de Contingência Web:\n\n"
                 "• Nível 1: Playwright HTTP (http://localhost:8000/...)\n"
                 "  -> Execução padrão com scraping da tabela HTML.\n\n"
                 "• Nível 2: Playwright Local via file://\n"
                 "  -> Em caso de timeout/porta 8000 fechada, abre nova aba limpa no arquivo local.\n\n"
                 "• Nível 3: Fallback Requests + BeautifulSoup\n"
                 "  -> Ativado se houver falha de renderizador gráfico do Chromium.\n\n"
                 "• Nível 4: Leitura direta em disco (0% de parada na fábrica).",
                 Emu(2900000), Emu(550000), Emu(2500000), Emu(2150000))
            ],
            "notes": "NOTAS DO ORADOR (👤 INTEGRANTE 3):\n- Assumir a apresentação respondendo: 'Como o robô trata erros?'.\n- Explicar o caso do Fornecedor D (bloqueado e dados negativos) e como o Data Filtering Gateway evita distorcer as notas dos demais.\n- Demonstrar a resiliência de rede em 4 camadas."
        },

        # SLIDE 8: Testes Automatizados Pytest (👤 Integrante 3)
        {
            "titulo": "Suíte de Testes Automatizados com Pytest (51 Testes)",
            "integrante": "Integrante 3 (Qualidade / MCDA)",
            "cards": [
                ("ESTRUTURA DA SUÍTE DE TESTES",
                 "🧪 test_coleta.py (24 testes):\n"
                 "• Varredura de diretórios e filtros de arquivos temporários (~$*, .*).\n"
                 "• Playwright RPA com Chromium real e mocks de erro de rede.\n"
                 "• Transição de abas no Playwright durante falhas de porta.\n\n"
                 "🧪 test_leitura.py (25 testes):\n"
                 "• Ingestão de XLSX e CSV (delimitadores ';' e ',').\n"
                 "• Tolerância a encodings (utf-8, latin-1) e acentuações.\n"
                 "• Leitura de critérios e template com fallback inteligente.\n\n"
                 "🧪 test_integration.py (2 testes):\n"
                 "• Teste ponta a ponta Coleta + Leitura + Main.",
                 Emu(350000), Emu(550000), Emu(2450000), Emu(2150000)),
                ("COMANDO & RESULTADOS DE COBERTURA",
                 "💻 Comando de Execução:\n"
                 "   python -m pytest -v --cov=src\n\n"
                 "📊 Métricas de Execução:\n"
                 "• Total de Testes: 51\n"
                 "• Status: 51 Aprovados (100% Passed ✅)\n"
                 "• Tempo de Execução: ~35 segundos\n\n"
                 "🎯 Cobertura de Código:\n"
                 "• src/etapa1_coleta/:  100% de cobertura\n"
                 "• src/etapa2_leitura/: 100% de cobertura\n"
                 "• src/logger.py:       84% de cobertura\n\n"
                 "🛡️ Garantia de Não-Regressão em merges para develop e main.",
                 Emu(2900000), Emu(550000), Emu(2500000), Emu(2150000))
            ],
            "notes": "NOTAS DO ORADOR (👤 INTEGRANTE 3):\n- Apresentar a suíte com 51 testes automatizados cobrindo todos os ramos de negócio e exceção com 100% de sucesso.\n- Destacar a integração no CI/CD."
        },

        # SLIDE 9: Algoritmo MCDA Ponderado (👤 Integrante 3)
        {
            "titulo": "Motor de Decisão Multicritério: Algoritmo MCDA",
            "integrante": "Integrante 3 (Qualidade / MCDA)",
            "cards": [
                ("PESOS & CRITÉRIOS DE NEGÓCIO",
                 "📋 Matriz de Critérios (criterios_ranking.xlsx):\n\n"
                 "• Custo: Peso 40% (0.40) | Direção: Menor é Melhor 🔻\n"
                 "• Prazo: Peso 25% (0.25) | Direção: Menor é Melhor 🔻\n"
                 "• Capacidade: Peso 20% (0.20) | Direção: Maior é Melhor 🔺\n"
                 "• Qualidade: Peso 15% (0.15) | Direção: Maior é Melhor 🔺\n\n"
                 "⚖️ Validação da Ponderação:\n"
                 "   Soma dos Pesos = 0.40 + 0.25 + 0.20 + 0.15 = 1.00 (100%)\n\n"
                 "🔄 Fallback de Negócio: Se a planilha for excluída, o robô aplica as constantes oficiais DEFAULT_WEIGHTS.",
                 Emu(350000), Emu(550000), Emu(2450000), Emu(2150000)),
                ("FÓRMULAS DE NORMALIZAÇÃO RELATIVA",
                 "📐 Normalização Relativa ao Benchmark de Mercado:\n\n"
                 "1. Critérios de Menor Valor (Custo e Prazo):\n"
                 "   Score = (Menor Valor do Mercado / Valor da Proposta) * 100\n"
                 "   -> A melhor cotação recebe nota 100.\n\n"
                 "2. Critérios de Maior Valor (Capacidade e Qualidade):\n"
                 "   Score = (Valor da Proposta / Maior Valor do Mercado) * 100\n"
                 "   -> A maior capacidade/qualidade recebe nota 100.\n\n"
                 "🎯 Equação do Score Global Ponderado:\n"
                 "   Nota Final = (Custo*0.40) + (Prazo*0.25) + (Capac*0.20) + (Qual*0.15)",
                 Emu(2900000), Emu(550000), Emu(2500000), Emu(2150000))
            ],
            "notes": "NOTAS DO ORADOR (👤 INTEGRANTE 3):\n- Explicar as fórmulas matemáticas do MCDA e o benchmark de mercado.\n- Passar a palavra para o Integrante 4."
        },

        # SLIDE 10: Docker & Docker Compose (👤 Integrante 4)
        {
            "titulo": "Containerização com Docker & Docker Compose",
            "integrante": "Integrante 4 (DevOps / Governança)",
            "cards": [
                ("DOCKERFILE & REQUISITOS EDITAL",
                 "🐳 Imagem Base Python 3.12-slim:\n"
                 "• Otimizada e leve para produção industrial.\n\n"
                 "🕒 Timezone Oficial: TZ=America/Manaus\n"
                 "• Configurado no sistema operacional (tzdata) e nas variáveis de ambiente.\n"
                 "• Garante que os carimbos de auditoria SOX sigam o horário da fábrica de Manaus.\n\n"
                 "🌐 Chromium Playwright Embutido:\n"
                 "• 'playwright install --with-deps chromium' na construção da imagem.\n\n"
                 "📂 Volumes Persistentes:\n"
                 "• ./output -> /app/output (Planilha de ranking)\n"
                 "• ./logs   -> /app/logs   (Auditoria JSON e logs)",
                 Emu(350000), Emu(550000), Emu(2450000), Emu(2150000)),
                ("ORQUESTRAÇÃO COM DOCKER COMPOSE",
                 "🚀 Execução Completa com 1 Único Comando:\n"
                 "   docker compose up --build\n\n"
                 "🛠️ Serviços Orquestrados:\n"
                 "1. Serviço 'web-panel':\n"
                 "   • Sobe servidor HTTP na porta 8000 simulando o portal corporativo da LG.\n"
                 "   • Healthcheck automático validando a disponibilidade do endpoint.\n\n"
                 "2. Serviço 'robot':\n"
                 "   • Depende do 'web-panel' estar ativo ('depends_on').\n"
                 "   • Executa a esteira ponta a ponta (ETL, MCDA e exportação).\n\n"
                 "🔒 Isolamento total sem dependências locais no host.",
                 Emu(2900000), Emu(550000), Emu(2500000), Emu(2150000))
            ],
            "notes": "NOTAS DO ORADOR (👤 INTEGRANTE 4):\n- Assumir a apresentação abordando a infraestrutura de Docker e Docker Compose.\n- Explicar os volumes persistentes, as variáveis de ambiente e o timezone America/Manaus."
        },

        # SLIDE 11: CI/CD GitHub Actions & GHCR (👤 Integrante 4)
        {
            "titulo": "Pipeline de CI/CD (GitHub Actions) & Publicação GHCR",
            "integrante": "Integrante 4 (DevOps / Governança)",
            "cards": [
                ("PIPELINE GITHUB ACTIONS (.github/workflows/ci.yml)",
                 "🔄 Gatilhos Automáticos:\n"
                 "• Push e Pull Request em main, develop e feature/**.\n\n"
                 "⚙️ Job 1: Lint, Testes e Cobertura:\n"
                 "1. Checkout do código e Setup Python 3.12 com cache pip.\n"
                 "2. Instalação de dependências e Chromium do Playwright.\n"
                 "3. Verificação de Linting estrito com Flake8 (0 erros).\n"
                 "4. Execução dos 51 testes Pytest com relatório de cobertura XML/Term.\n\n"
                 "⚙️ Job 2: Build e Validação Docker:\n"
                 "• Validação da integridade da imagem Docker após testes.",
                 Emu(350000), Emu(550000), Emu(2450000), Emu(2150000)),
                ("PUBLICAÇÃO & EXECUÇÃO VIA GHCR",
                 "📦 GitHub Container Registry (GHCR):\n"
                 "• Imagem homologada publicada em ghcr.io.\n"
                 "• Não utilizamos Docker Hub, atendendo estritamente ao edital.\n\n"
                 "💻 Comandos de Demonstração em Produção:\n"
                 "1. Baixar a imagem publicada:\n"
                 "   docker pull ghcr.io/sannyer3232/hyperautomation-av3-equipe1:latest\n\n"
                 "2. Executar o container com mapeamento de volumes:\n"
                 "   docker run --rm -v $(pwd)/output:/app/output \\\n"
                 "              -v $(pwd)/logs:/app/logs \\\n"
                 "              -e TZ=America/Manaus \\\n"
                 "              ghcr.io/sannyer3232/hyperautomation-av3-equipe1:latest",
                 Emu(2900000), Emu(550000), Emu(2500000), Emu(2150000))
            ],
            "notes": "NOTAS DO ORADOR (👤 INTEGRANTE 4):\n- Explicar a esteira automatizada de CI/CD no GitHub Actions e a publicação no GitHub Container Registry (GHCR).\n- Mostrar os comandos docker pull e docker run."
        },

        # SLIDE 12: Logs e Auditoria SOX (👤 Integrante 4)
        {
            "titulo": "Rastreabilidade Digital, Logs e Auditoria SOX",
            "integrante": "Integrante 4 (DevOps / Governança)",
            "cards": [
                ("TRILHA DE AUDITORIA (logs/auditoria.json)",
                 "🏛️ Snapshot Estruturado para Compliance SOX:\n\n"
                 "• Data/Hora de Início e Fim (ISO 8601 com TZ Manaus).\n"
                 "• Total de propostas recebidas (4) e arquivos lidos.\n"
                 "• Lista de propostas aprovadas com dados comerciais.\n"
                 "• Justificativa da desclassificação do Fornecedor D.\n"
                 "• Memória de cálculo detalhada com scores parciais (Custo, Prazo, Capacidade, Qualidade).\n"
                 "• Ranking final homologado com classificação oficial.\n"
                 "• Registro de erros e exceções capturadas.",
                 Emu(350000), Emu(550000), Emu(2450000), Emu(2150000)),
                ("LOGS OPERACIONAIS & ARTEFATO FINAL",
                 "📋 Log Operacional Contínuo (logs/execucao.log):\n"
                 "• Formato padronizado: %(asctime)s [%(levelname)s] %(message)s\n"
                 "• 4 Níveis de Severidade Utilizados:\n"
                 "  - [INFO]: Etapas normais, leitura e ranking gerado.\n"
                 "  - [WARNING]: Fallback ativado, Fornecedor D rejeitado.\n"
                 "  - [ERROR]: Falha técnica localizada na leitura.\n"
                 "  - [CRITICAL]: Erro fatal não tratado na pipeline.\n\n"
                 "📊 Planilha Final Homologada (output/ranking_final.xlsx):\n"
                 "• Preenchimento fiel do template modelo_ranking.xlsx.",
                 Emu(2900000), Emu(550000), Emu(2500000), Emu(2150000))
            ],
            "notes": "NOTAS DO ORADOR (👤 INTEGRANTE 4):\n- Explicar a diferença entre o log operacional para suporte e o snapshot JSON de auditoria SOX para compliance da LG."
        },

        # SLIDE 13: Conclusão e Valor Entregue (👤 Integrante 4)
        {
            "titulo": "Conclusão e Entrega de Valor para o Negócio",
            "integrante": "Integrante 4 (DevOps / Governança)",
            "cards": [
                ("O QUE DESENVOLVEMOS?",
                 "• Solução completa de Hyperautomation.\n"
                 "• RPA Playwright + Motor MCDA + Pytest + Docker + CI/CD + GHCR.\n"
                 "• 100% dos requisitos do edital da LG atendidos.",
                 Emu(350000), Emu(550000), Emu(1600000), Emu(2150000)),
                ("QUAL PROBLEMA RESOLVEU?",
                 "• Fim da digitação manual de propostas.\n"
                 "• Barramento automático de fornecedores irregulares.\n"
                 "• Redução do ciclo de compras de 8 horas para 25 segundos.",
                 Emu(2030000), Emu(550000), Emu(1650000), Emu(2150000)),
                ("QUAL VALOR ENTREGA?",
                 "• Decisões 100% objetivas e matemáticas.\n"
                 "• Total conformidade com auditorias SOX da LG.\n"
                 "• Escalabilidade para mais de 350 cotações/mês.",
                 Emu(3760000), Emu(550000), Emu(1650000), Emu(2150000))
            ],
            "notes": "NOTAS DO ORADOR (👤 INTEGRANTE 4):\n- Fazer o fechamento executivo respondendo às 5 perguntas obrigatórias de conclusão.\n- Anunciar o início da demonstração prática ao vivo."
        },

        # SLIDE 14: Roteiro da Demonstração Prática (👥 4 Integrantes)
        {
            "titulo": "Roteiro da Demonstração Prática (7 Etapas Obrigatórias)",
            "integrante": "4 Integrantes (Demonstração)",
            "cards": [
                ("ETAPAS 1 A 4 DA DEMONSTRAÇÃO",
                 "👤 Integrante 1 — ETAPA 1: Código & Repositório\n"
                 "• Mostrar pastas (src, tests, resources, docs, .github) e configurações.\n\n"
                 "👤 Integrante 2 — ETAPA 2: Execução da Pipeline\n"
                 "• Executar 'python src/main.py' no terminal mostrando o robô ao vivo.\n\n"
                 "👤 Integrante 3 — ETAPA 3: Testes Automatizados\n"
                 "• Executar 'python -m pytest -v' demonstrando 51 testes passando.\n\n"
                 "👤 Integrante 4 — ETAPA 4: Workflow no GitHub Actions\n"
                 "• Abrir navegador mostrando a esteira verde (Lint + Test + Docker).",
                 Emu(350000), Emu(550000), Emu(2450000), Emu(2150000)),
                ("ETAPAS 5 A 7 DA DEMONSTRAÇÃO",
                 "👤 Integrante 4 — ETAPA 5: Imagem no GHCR\n"
                 "• Mostrar o pacote publicado no GitHub Container Registry (ghcr.io).\n\n"
                 "👤 Integrante 4 — ETAPA 6: Execução via Docker Compose\n"
                 "• Demonstrar 'docker compose up' com timezone America/Manaus.\n\n"
                 "👥 Todos — ETAPA 7: Resultados, Logs & Auditoria\n"
                 "• Abrir 'output/ranking_final.xlsx' (1º B, 2º A, 3º C, Desclassificado D).\n"
                 "• Abrir 'logs/auditoria.json' mostrando a memória de cálculo e trilha SOX.\n"
                 "• Abrir 'logs/execucao.log' demonstrando severidades INFO e WARNING.",
                 Emu(2900000), Emu(550000), Emu(2500000), Emu(2150000))
            ],
            "notes": "NOTAS DO ORADOR (👥 TODOS OS 4 INTEGRANTES):\n- Seguir exatamente este roteiro cronometrado na demonstração ao vivo para a banca examinadora."
        }
    ]

    # Criação dos slides 2 a 14 a partir do slide de referência do template
    for idx_slide, dado in enumerate(dados_slides, start=2):
        s_novo = clonar_slide_template(prs, slide_conteudo_ref_idx)
        limpar_formas_conteudo_antigo(s_novo)
        ajustar_cabecalho_e_rodape(s_novo, dado["titulo"], f"‹#›/15", dado["integrante"])

        for card_info in dado["cards"]:
            tit_c, desc_c, cx, cy, cw, ch = card_info
            cor_bg = COLOR_CARD_BG

            c = adicionar_container_card(s_novo, cx, cy, cw, ch, tit_c, bg_color=cor_bg)
            tf = c.text_frame
            p = tf.add_paragraph()
            p.text = desc_c
            p.font.size = Pt(7.5)
            p.font.color.rgb = COLOR_DARK_TEXT

        s_novo.notes_slide.notes_text_frame.text = dado["notes"]

    # =========================================================================
    # SLIDE 15: SLIDE DE ENCERRAMENTO E PERGUNTAS (Template Slide 9)
    # =========================================================================
    s15 = clonar_slide_template(prs, slide_fim_ref_idx)
    for shape in s15.shapes:
        if shape.has_text_frame:
            texto = shape.text_frame.text
            if "Obrigado" in texto or "Dúvidas" in texto:
                shape.text_frame.clear()
                p = shape.text_frame.paragraphs[0]
                p.text = "Obrigado!\nDúvidas da Banca?"
                p.font.size = Pt(22)
                p.font.bold = True
                p.font.color.rgb = COLOR_WHITE
            elif "carvalho" in texto or "@" in texto:
                shape.text_frame.clear()
                p = shape.text_frame.paragraphs[0]
                p.text = "LG Electronics do Brasil • Equipe 01 • Turma T02 (2026)\n" \
                         "Repositório: github.com/Sannyer3232/hyperautomation-av3-equipe1\n" \
                         "Imagem GHCR: ghcr.io/sannyer3232/hyperautomation-av3-equipe1:latest"
                p.font.size = Pt(8.5)
                p.font.color.rgb = COLOR_CYAN

    s15.notes_slide.notes_text_frame.text = (
        "NOTAS DO ORADOR (👥 TODOS OS 4 INTEGRANTES):\n"
        "- Agradecer a atenção do professor Moisés Levy e da banca.\n"
        "- Abrir para a sessão de perguntas individuais e coletivas da banca examinadora."
    )

    # Deleta os slides originais não utilizados (do índice 1 ao 8 do template original)
    # Excluímos 8 vezes o índice 1 para que o Slide 1 (Capa editada) permaneça no índice 0
    # e os novos slides fiquem logo após.
    for _ in range(8):
        slide_id = prs.slides._sldIdLst[1]
        prs.slides._sldIdLst.remove(slide_id)

    caminho_final = Path("docs") / "Apresentacao_Defesa_Hyperautomation_LG_Equipe01.pptx"
    prs.save(str(caminho_final))
    print(f"Sucesso! Apresentação gerada diretamente no template oficial em: {caminho_final}")
    print(f"Total de slides na apresentação final: {len(prs.slides)}")


if __name__ == "__main__":
    gerar_apresentacao_oficial()
